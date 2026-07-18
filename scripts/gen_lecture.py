# -*- coding: utf-8 -*-
"""
3-hour lecture deck: "찰칵레시피 개발 여정 — 클로드 코드로 배우는 바이브 코딩"
Built with python-pptx, no network/template dependency.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

PROJECT_ROOT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..")
DIAGRAMS = os.path.join(PROJECT_ROOT, "docs", "diagrams")

FONT = "맑은 고딕"

CREAM = RGBColor(0xFB, 0xF4, 0xEC)
CREAM2 = RGBColor(0xF3, 0xE9, 0xDA)
INK = RGBColor(0x2C, 0x23, 0x1D)
MUTED = RGBColor(0x8A, 0x7B, 0x6E)
CORAL = RGBColor(0xD9, 0x64, 0x3A)
CORAL_DARK = RGBColor(0xB1, 0x4B, 0x26)
LINE = RGBColor(0xE8, 0xDC, 0xCB)
GREEN = RGBColor(0x3E, 0x81, 0x56)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x2E, 0x3A, 0x4E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PAGE = {"n": 0}


def new_slide(prs, bg=CREAM):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()
    bg_shape.shadow.inherit = False
    # send to back
    sp = bg_shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    PAGE["n"] += 1
    return slide


def _set_text(tf, text, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT,
              font=FONT, line_spacing=1.15, italic=False):
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color


def add_text(slide, text, left, top, width, height, size=20, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.15,
             italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    _set_text(tf, text, size=size, color=color, bold=bold, align=align, font=font,
              line_spacing=line_spacing, italic=italic)
    return box


def add_bullets(slide, items, left, top, width, height, size=18, color=INK, font=FONT,
                 space_after=10, line_spacing=1.2):
    """items: list of (text, level, opts) where opts is dict with optional
    'bold','color','bullet' keys. level 0/1/2 controls indent."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            text, level, opts = item, 0, {}
        elif len(item) == 2:
            text, level = item
            opts = {}
        else:
            text, level, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        indent = level * 0.35
        bullet_char = opts.get("bullet", ["●", "–", "·"][min(level, 2)])
        prefix = f"{bullet_char}  " if bullet_char else ""
        r = p.add_run()
        r.text = ("    " * level) + prefix + text
        r.font.size = Pt(opts.get("size", size - level * 1.5))
        r.font.bold = opts.get("bold", level == 0)
        r.font.name = font
        r.font.color.rgb = opts.get("color", color if level == 0 else MUTED)
    return box


def add_rect(slide, left, top, width, height, fill=None, line_color=None, line_w=1.0,
             radius=None, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if radius:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = shadow
    return shp


def add_image_fit(slide, path, left, top, max_w, max_h, align="center"):
    im = Image.open(path)
    ratio = im.width / im.height
    w = max_w
    h = w / ratio
    if h > max_h:
        h = max_h
        w = h * ratio
    if align == "center":
        left = left + (max_w - w) / 2
    pic = slide.shapes.add_picture(path, left, top + (max_h - h) / 2, width=int(w), height=int(h))
    return pic


def footer(slide, label, page_no, total):
    add_text(slide, label, Inches(0.6), Inches(7.08), Inches(8), Inches(0.35),
              size=10, color=MUTED, font=FONT)
    add_text(slide, f"{page_no} / {total}", Inches(12.2), Inches(7.08), Inches(0.9), Inches(0.35),
              size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT)


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def section_divider(prs, no, title, subtitle, minutes):
    slide = new_slide(prs, bg=CORAL_DARK)
    add_text(slide, f"PART {no}", Inches(0.9), Inches(2.3), Inches(4), Inches(0.6),
             size=20, color=CREAM2, bold=True, font=FONT)
    add_text(slide, title, Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.6),
             size=40, color=WHITE, bold=True, font=FONT)
    add_text(slide, subtitle, Inches(0.9), Inches(4.1), Inches(11), Inches(1.0),
             size=18, color=CREAM2, font=FONT)
    add_rect(slide, Inches(0.9), Inches(4.85), Inches(1.4), Inches(0.45), fill=WHITE, radius=0.5)
    add_text(slide, f"약 {minutes}분", Inches(0.9), Inches(4.85), Inches(1.4), Inches(0.45),
              size=14, color=CORAL_DARK, bold=True, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE, font=FONT)
    return slide


def content_slide(prs, kicker, title, page_no, total, bg=CREAM):
    slide = new_slide(prs, bg=bg)
    if kicker:
        add_text(slide, kicker, Inches(0.7), Inches(0.42), Inches(10), Inches(0.4),
                  size=14, color=CORAL_DARK, bold=True, font=FONT)
    add_text(slide, title, Inches(0.7), Inches(0.75), Inches(11.9), Inches(1.0),
              size=30, color=INK, bold=True, font=FONT)
    add_rect(slide, Inches(0.7), Inches(1.62), Inches(1.1), Pt(3), fill=CORAL)
    footer(slide, "찰칵레시피로 배우는 바이브 코딩", page_no, total)
    return slide


def site_slide(prs, page_no, total, tag, domain, role_kr, analogy, bullets, notes, accent):
    slide = content_slide(prs, tag, domain, page_no, total)
    add_rect(slide, Inches(0.7), Inches(1.95), Inches(4.5), Inches(0.62), fill=accent, radius=0.5)
    add_text(slide, role_kr, Inches(0.7), Inches(1.95), Inches(4.5), Inches(0.62),
              size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, analogy, Inches(0.7), Inches(2.75), Inches(11.6), Inches(0.9),
              size=17, color=MUTED, italic=True)
    add_bullets(slide, bullets, Inches(0.7), Inches(3.7), Inches(11.8), Inches(3.1), size=18)
    set_notes(slide, notes)
    return slide


def competency_slide(prs, page_no, total, idx, title_kr, essence, bullets, notes):
    slide = new_slide(prs, bg=CORAL_DARK if idx % 2 else NAVY)
    add_text(slide, f"{idx:02d}", Inches(0.7), Inches(0.5), Inches(3), Inches(1.6),
              size=64, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    add_text(slide, title_kr, Inches(0.7), Inches(2.0), Inches(11.8), Inches(1.1),
              size=32, color=WHITE, bold=True)
    add_text(slide, essence, Inches(0.7), Inches(3.05), Inches(11.6), Inches(0.8),
              size=18, color=CREAM2, italic=True)
    add_bullets(slide, [(t, 0, {"color": CREAM2, "bold": False}) for t in bullets],
                Inches(0.7), Inches(3.9), Inches(11.8), Inches(3.0), size=18, color=CREAM2)
    footer(slide, "AI 활용 개발자에게 필요한 7가지 역량", page_no, total)
    set_notes(slide, notes)
    return slide


def process_slide(prs, page_no, total, step_no, title, desc, bullets, notes, image=None):
    slide = content_slide(prs, f"개발 과정 {step_no}단계", title, page_no, total)
    body_w = Inches(6.6) if image else Inches(11.8)
    add_text(slide, desc, Inches(0.7), Inches(1.95), body_w, Inches(0.9), size=17, color=MUTED, italic=True)
    add_bullets(slide, bullets, Inches(0.7), Inches(2.85), body_w, Inches(3.9), size=17)
    if image:
        add_image_fit(slide, image, Inches(7.5), Inches(1.95), Inches(5.15), Inches(4.6))
    footer(slide, "찰칵레시피로 배우는 바이브 코딩", page_no, total)
    set_notes(slide, notes)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    TOTAL = 50

    # 1. Title
    s = new_slide(prs, bg=CREAM)
    add_rect(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill=CORAL)
    add_text(s, "대학 신입생을 위한 3시간 특강", Inches(0.9), Inches(1.6), Inches(10), Inches(0.5),
              size=18, color=CORAL_DARK, bold=True)
    add_text(s, "찰칵레시피 개발 여정", Inches(0.9), Inches(2.15), Inches(11.5), Inches(1.3),
              size=48, color=INK, bold=True)
    add_text(s, "클로드 코드로 배우는 바이브 코딩\n— 아이디어에서 실제 배포까지", Inches(0.9), Inches(3.35),
              Inches(11.5), Inches(1.3), size=24, color=INK, line_spacing=1.3)
    add_text(s, "냉장고 사진 한 장 → 오늘 뭐 해먹지, 그 답을 만든 이야기", Inches(0.9), Inches(4.7),
              Inches(11), Inches(0.6), size=16, color=MUTED, italic=True)
    add_rect(s, Inches(0.9), Inches(6.3), Inches(11.5), Pt(1.2), fill=LINE)
    add_text(s, "40년차 개발자가 신입생에게 건네는 첫 코딩 수업", Inches(0.9), Inches(6.5), Inches(11.5),
              Inches(0.5), size=14, color=MUTED)
    set_notes(s, "환영 인사. 오늘 3시간 동안 코드를 한 줄도 직접 치지 않고도, 실제로 인터넷에 배포된 "
                 "레시피 추천 앱을 어떻게 만들었는지 처음부터 끝까지 함께 따라가 볼 것이라고 소개한다. "
                 "이 앱은 지어낸 예시가 아니라 이 강의를 준비하며 실제로 만들고 배포한 서비스임을 강조한다.")

    n = 1

    def nx():
        nonlocal n
        n += 1
        return n

    # 2. 강의 소개
    s = content_slide(prs, "여는 말", "오늘 3시간, 우리가 함께 할 일", nx(), TOTAL)
    add_bullets(s, [
        ("냉장고 사진 한 장으로 시작해 실제 배포된 웹앱이 완성되기까지, 그 전 과정을 그대로 따라간다", 0),
        ("코드 문법을 외우는 수업이 아니다 — \"AI에게 무엇을, 어떻게 요청하는가\"를 배우는 수업이다", 0),
        ("이 강의가 끝나면 여러분 손에 남는 것", 0),
        ("① 실제 서비스 하나가 어떤 부품들로 이루어지는지 보는 눈", 1, {"bullet": "→"}),
        ("② 클로드 코드와 대화하며 만드는 \"바이브 코딩\"의 감각", 1, {"bullet": "→"}),
        ("③ \"나도 만들 수 있겠다\"는 근거 있는 자신감", 1, {"bullet": "→"}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "이번 슬라이드는 강의 전체의 지도 역할을 한다. 신입생들에게 오늘은 문법 암기가 아니라 "
                 "실제 제품이 만들어지는 과정을 처음부터 끝까지 보여주는 시간임을 명확히 한다. "
                 "특히 세 가지 목표(보는 눈, 바이브 코딩 감각, 자신감)를 칠판이나 슬라이드에 남겨두고 "
                 "강의 말미에 다시 이 슬라이드로 돌아와 확인해도 좋다.")

    # 3. 강사의 철학
    s = content_slide(prs, "시작하기 전에", "40년을 코딩했지만, 요즘처럼 빠른 적은 없었습니다", nx(), TOTAL)
    add_bullets(s, [
        ("예전에는 화면 하나 만드는 데도 문법책을 뒤적이며 며칠이 걸렸습니다", 0),
        ("지금은 \"이런 앱을 만들고 싶다\"고 말하면, AI가 그 자리에서 뼈대를 만들어 보여줍니다", 0),
        ("그렇다고 아무나 아무거나 뚝딱 만들어내는 건 아닙니다 — 무엇을 원하는지 \"정확히 설명하는 능력\"이 "
         "새로운 핵심 실력이 되었습니다", 0),
        ("이 강의는 그 능력을, 우리가 실제로 만든 앱을 교재 삼아 함께 연습하는 자리입니다", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "이 슬라이드에서 진심을 담아 이야기한다. 40년의 경력이 주는 신뢰감을 바탕으로 "
                 "\"기술이 이렇게 빨리 바뀐 적이 없었다\"는 것을 솔직하게 인정하면서도, "
                 "AI가 모든 걸 대신해주는 것이 아니라 \"정확히 설명하는 능력\"이 여전히, "
                 "아니 오히려 더 중요해졌다는 메시지로 이어간다. 이 메시지가 강의 마지막 파트(7가지 역량)의 "
                 "복선이 된다는 것을 은근히 암시해도 좋다.")

    # 4. 시간 배분
    s = content_slide(prs, "오늘의 순서", "3시간, 6개의 파트로 나눠 갑니다", nx(), TOTAL)
    rows = [
        ("Part 1", "우리가 만든 것 — 찰칵레시피", "20분"),
        ("Part 2", "바이브 코딩이란 무엇인가", "25분"),
        ("쉬는 시간", "", "10분"),
        ("Part 3", "기술 스택 이해하기 — 4개 사이트의 역할", "40분"),
        ("Part 4", "개발 과정 실제로 따라가보기", "35분"),
        ("쉬는 시간", "", "10분"),
        ("Part 5", "여러분도 할 수 있다는 확신", "20분"),
        ("Part 6", "AI 활용 개발자에게 필요한 7가지 역량", "25분"),
    ]
    top = Inches(2.0)
    row_h = Inches(0.52)
    for i, (a, b, c) in enumerate(rows):
        y = top + row_h * i
        is_break = a == "쉬는 시간"
        bg = CREAM2 if not is_break else RGBColor(0xEF, 0xE3, 0xD3)
        add_rect(s, Inches(0.7), y, Inches(11.8), row_h - Pt(4), fill=bg, radius=0.15)
        add_text(s, a, Inches(0.95), y, Inches(1.9), row_h - Pt(4), size=15,
                  color=CORAL_DARK if not is_break else MUTED, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, b, Inches(2.9), y, Inches(7.7), row_h - Pt(4), size=15, color=INK,
                  anchor=MSO_ANCHOR.MIDDLE, italic=is_break)
        add_text(s, c, Inches(10.7), y, Inches(1.6), row_h - Pt(4), size=15, color=MUTED, bold=True,
                  align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    set_notes(s, "전체 진행표를 미리 보여주면서 언제 쉬는지 안내해 학생들의 긴장을 풀어준다. "
                 "표 그대로 진행하되, 현장 반응에 따라 Part별 시간은 ±5분 정도 유연하게 조정해도 된다고 "
                 "말해준다.")

    # ---------------- PART 1 ----------------
    section_divider(prs, 1, "우리가 만든 것", "찰칵레시피 — 냉장고 사진 한 장의 힘", 20)
    nx()

    s = content_slide(prs, "Part 1 · 우리가 만든 것", "찰칵레시피는 어떤 문제를 풀어주는 앱인가", nx(), TOTAL)
    add_bullets(s, [
        ("모두가 겪는 아주 사소하지만 매일 반복되는 고민 — \"오늘 뭐 해먹지?\"", 0),
        ("냉장고 문을 열어보면 재료는 있는데, 그걸로 뭘 만들 수 있는지는 막상 떠오르지 않는다", 0),
        ("찰칵레시피의 답: 냉장고 사진을 찰칵 찍어 올리면", 0, {"bold": True, "color": CORAL_DARK}),
        ("AI가 사진 속 재료를 알아보고 (비전 AI)", 1, {"bullet": "①"}),
        ("그 재료로 만들 수 있는 레시피 3개를 추천해준다 (텍스트 AI)", 1, {"bullet": "②"}),
        ("특별한 기능이 아니라, \"누구나 공감하는 작은 불편\"에서 출발했다는 점이 중요하다", 0,
         {"bold": True}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "좋은 프로젝트는 거창한 아이디어에서 시작하지 않는다는 것을 강조. 학생들에게 "
                 "\"여러분이 오늘 아침 겪은 사소한 불편은 무엇이었나요?\"라고 질문을 던져 잠깐 대화를 "
                 "유도해도 좋다. 이 질문이 Part 5의 \"무엇이든 만들 수 있다\" 파트와 자연스럽게 연결된다.")

    s = content_slide(prs, "Part 1 · 우리가 만든 것", "사용자 흐름 4단계", nx(), TOTAL)
    steps = [
        ("① 사진 업로드", "냉장고 안을 최대 3장까지 찍어 올린다"),
        ("② 재료 확인", "AI가 알아낸 재료 목록을 사람이 직접 고치고 보완한다"),
        ("③ 레시피 추천", "확정된 재료 + 내 취향(맵기/난이도/시간)으로 레시피 3개를 받는다"),
        ("④ 평가와 저장", "마음에 든 레시피는 저장하고, 좋아요/코멘트로 기록을 남긴다"),
    ]
    box_w = Inches(2.78)
    gap = Inches(0.15)
    for i, (title, desc) in enumerate(steps):
        x = Inches(0.7) + (box_w + gap) * i
        add_rect(s, x, Inches(2.15), box_w, Inches(3.7), fill=WHITE, line_color=LINE, line_w=1.2, radius=0.08)
        add_rect(s, x, Inches(2.15), box_w, Inches(0.85), fill=CORAL, radius=0.08)
        add_text(s, title, x, Inches(2.15), box_w, Inches(0.85), size=17, color=WHITE, bold=True,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc, x + Inches(0.18), Inches(3.15), box_w - Inches(0.36), Inches(2.5),
                  size=14.5, color=INK, line_spacing=1.3)
        if i < 3:
            add_text(s, "→", x + box_w, Inches(2.15), gap + Inches(0.3), Inches(0.85), size=20,
                      color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    set_notes(s, "화면 상단에 4단계 진행 표시(마법사 스텝)가 실제로 붙어 있다는 점을 언급하면 좋다. "
                 "\"이 4단계가 곧 오늘 우리가 만드는 전체 프로그램의 뼈대\"라는 것을 짚어준다.")

    s = content_slide(prs, "Part 1 · 우리가 만든 것", "왜 이 프로젝트가 좋은 교재인가", nx(), TOTAL)
    add_bullets(s, [
        ("실제로 배포되어 지금 이 순간에도 접속할 수 있는 살아있는 서비스다 (recipe4fridge-pic.vercel.app)", 0),
        ("회원가입/로그인, 데이터베이스, 외부 AI API, 배포까지 — 실무 웹서비스의 핵심 요소를 모두 담고 있다", 0),
        ("그러면서도 복잡한 결제·물류 같은 요소는 없어 신입생이 전체 그림을 한 번에 이해하기 좋다", 0),
        ("무엇보다, 오늘 보여드릴 것은 \"완성된 결과\"가 아니라 \"만들어져 온 과정\" 그 자체다", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "이 슬라이드는 Part1의 마무리이자 Part2로 넘어가는 다리 역할. \"그럼 이걸 어떻게 만들었는가\"로 "
                 "자연스럽게 화제를 전환한다.")

    # ---------------- PART 2 ----------------
    section_divider(prs, 2, "바이브 코딩이란 무엇인가", "AI와 대화하며 만드는 새로운 개발 방식", 25)
    nx()

    s = content_slide(prs, "Part 2 · 바이브 코딩", "전통적 개발 vs 바이브 코딩", nx(), TOTAL)
    col_w = Inches(5.7)
    add_rect(s, Inches(0.7), Inches(2.0), col_w, Inches(4.2), fill=WHITE, line_color=LINE, line_w=1.2, radius=0.06)
    add_text(s, "전통적인 개발", Inches(0.7), Inches(2.15), col_w, Inches(0.5), size=19, bold=True,
              color=INK, align=PP_ALIGN.CENTER)
    add_bullets(s, [
        ("문법을 하나하나 배우고 외운다", 0),
        ("에러 메시지를 검색하며 원인을 추적한다", 0),
        ("기능 하나 추가에도 관련 문서를 오래 뒤진다", 0),
        ("\"어떻게 구현하는가\"에 대부분의 시간을 쓴다", 0),
    ], Inches(0.95), Inches(2.75), col_w - Inches(0.5), Inches(3.2), size=17)
    add_rect(s, Inches(6.65), Inches(2.0), col_w, Inches(4.2), fill=CORAL, radius=0.06)
    add_text(s, "바이브 코딩", Inches(6.65), Inches(2.15), col_w, Inches(0.5), size=19, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
    add_bullets(s, [
        ("무엇을 원하는지 사람의 말로 설명한다", 0, {"color": CREAM2}),
        ("AI가 코드를 직접 작성·실행하고 결과를 보여준다", 0, {"color": CREAM2}),
        ("사람은 결과를 보고 \"이건 맞다/아니다\"를 판단한다", 0, {"color": CREAM2}),
        ("\"무엇을, 왜 만드는가\"에 대부분의 시간을 쓴다", 0, {"color": CREAM2}),
    ], Inches(6.9), Inches(2.75), col_w - Inches(0.5), Inches(3.2), size=17, color=CREAM2)
    set_notes(s, "두 방식이 대립한다기보다, 힘을 쏟는 지점이 옮겨갔다는 것을 강조. 문법은 AI가 대신 "
                 "처리해주는 대신, \"무엇을 왜 만드는가\"를 설명하는 능력이 사람에게 더 중요해졌다는 메시지.")

    s = content_slide(prs, "Part 2 · 바이브 코딩", "클로드 코드(Claude Code)란?", nx(), TOTAL)
    add_bullets(s, [
        ("터미널(명령창)에서 실행하는 AI 개발 동료", 0, {"bold": True}),
        ("코드를 \"읽고 → 쓰고 → 실행하고 → 결과를 확인\"하는 것을 스스로 반복할 수 있다", 0),
        ("파일을 만들고 고치고, 터미널 명령을 실행하고, 브라우저까지 직접 열어 화면을 확인한다", 0),
        ("사람은 \"방향을 정하고 결과를 검수하는\" 역할에 집중할 수 있다", 0),
        ("오늘 이 강의 자료도, 찰칵레시피 앱도 전부 클로드 코드와의 대화로 만들어졌다", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "클로드 코드를 단순 \"챗봇\"이 아니라 \"직접 일하는 동료\"로 소개하는 것이 핵심. "
                 "코드를 읽고 실행 결과를 스스로 확인한다는 점에서 기존 챗봇과 다르다는 것을 짚어준다.")

    s = content_slide(prs, "Part 2 · 바이브 코딩", "바이브 코딩의 4단계 루프", nx(), TOTAL)
    loop = [
        ("① 설명한다", "\"이런 화면/기능이 필요해\"라고 사람의 말로 요청"),
        ("② 구현한다", "AI가 코드를 작성하고 필요하면 직접 실행"),
        ("③ 검증한다", "화면을 함께 보며 \"이게 맞나?\" 확인"),
        ("④ 반복한다", "안 맞으면 다시 설명 — 이 사이클을 계속 돈다"),
    ]
    box_w = Inches(2.78)
    for i, (title, desc) in enumerate(loop):
        x = Inches(0.7) + (box_w + Inches(0.15)) * i
        add_rect(s, x, Inches(2.3), box_w, Inches(3.3), fill=CREAM2, radius=0.1)
        add_text(s, title, x, Inches(2.5), box_w, Inches(0.6), size=19, bold=True, color=CORAL_DARK,
                  align=PP_ALIGN.CENTER)
        add_text(s, desc, x + Inches(0.2), Inches(3.15), box_w - Inches(0.4), Inches(2.2), size=14.5,
                  color=INK, line_spacing=1.3, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(s, "↻" if i == 2 else "→", x + box_w, Inches(2.3), Inches(0.3), Inches(3.3),
                      size=20, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "④ 다음에는 다시 ①로 — 이 루프를 수백 번 돌린 결과가 바로 오늘의 찰칵레시피다",
              Inches(0.7), Inches(5.9), Inches(11.8), Inches(0.6), size=15, color=MUTED, italic=True,
              align=PP_ALIGN.CENTER)
    set_notes(s, "이 루프 다이어그램이 오늘 강의 전체를 관통하는 핵심 그림이다. Part4에서 실제 개발 "
                 "단계들을 설명할 때 이 루프가 몇 번이고 반복되었다는 걸 계속 상기시켜준다.")

    s = content_slide(prs, "Part 2 · 바이브 코딩", "오늘 실제로 나눴던 대화, 미리 살짝 보기", nx(), TOTAL)
    add_bullets(s, [
        ("\"관리자만 보던 레시피 평가를, 회원 누구나 볼 수 있게 바꿔줘\"", 0, {"bullet": "“"}),
        ("\"사진과 재료 인식 결과를 사용자와 AI가 함께 채점하게 만들고 싶어\"", 0, {"bullet": "“"}),
        ("\"모델 드롭다운 UI가 카드 밖으로 삐져나와\"", 0, {"bullet": "“"}),
        ("세 문장 모두 코드 용어가 하나도 없다 — 그런데 전부 실제로 구현됐다", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "실제로 이 강의를 준비하며 나눈 대화에서 그대로 가져온 문장들이라는 걸 밝힌다. "
                 "특히 \"코드 용어가 하나도 없다\"는 점을 강조해서 신입생들의 심리적 장벽을 낮춰준다.")

    # ---- BREAK 1 ----
    s = new_slide(prs, bg=NAVY)
    add_text(s, "잠시 쉬어갑니다", Inches(0.9), Inches(2.9), Inches(11), Inches(1.0), size=40, bold=True,
              color=WHITE)
    add_text(s, "10분 후, \"기술 스택 이해하기\"에서 다시 만나요", Inches(0.9), Inches(4.0), Inches(11),
              Inches(0.6), size=18, color=CREAM2)
    nx()

    # ---------------- PART 3 ----------------
    section_divider(prs, 3, "기술 스택 이해하기", "4개의 외부 서비스가 각자의 역할을 맡는다", 40)
    nx()

    s = content_slide(prs, "Part 3 · 기술 스택", "전체 아키텍처 한눈에 보기", nx(), TOTAL)
    add_image_fit(s, os.path.join(DIAGRAMS, "architecture.png"), Inches(0.7), Inches(1.95),
                  Inches(11.9), Inches(4.8))
    set_notes(s, "이 그림 한 장이 오늘 Part3의 지도가 된다. 사용자의 브라우저에서 시작해 Next.js 서버를 "
                 "거쳐 비전 API, 텍스트 API를 호출하고, Supabase에 저장되는 전체 흐름을 짚어준다. "
                 "이제부터 이 그림에 나오는 각 상자가 실제로 어떤 웹사이트인지 하나씩 열어본다고 예고한다.")

    s = content_slide(prs, "Part 3 · 기술 스택", "오케스트라 비유 — 나는 지휘자다", nx(), TOTAL)
    add_bullets(s, [
        ("혼자 모든 악기를 다 배워서 연주할 필요는 없다", 0),
        ("코드(Next.js)는 지휘자의 악보 — 누가 언제 연주할지 정리한다", 0),
        ("vercel.com은 공연장 — 관객(사용자)이 실제로 입장해서 공연을 볼 수 있는 곳", 0),
        ("supabase.com은 악보 보관소이자 대기실 — 회원 명단, 데이터, 사진을 보관", 0),
        ("openrouter.ai는 초청 연주자 섭외처 — 필요한 순간에 AI 연주자(모델)를 빌려온다", 0),
        ("console.cloud.google.com은 신분 확인 데스크 — \"이 관객은 진짜 맞다\"를 확인해주는 곳", 0),
        ("우리가 할 일은 이 4곳과 \"대화 잘 하는 지휘자\"가 되는 것이다", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.5), size=18)
    set_notes(s, "비유를 천천히, 여유 있게 설명한다. 신입생들이 \"4개나 배워야 하나\"라는 부담을 느끼지 "
                 "않도록, 각 서비스가 하는 일이 사실 우리가 일상에서 이미 이해하고 있는 개념(공연장, 창고, "
                 "섭외처, 신분증)의 연장선이라는 걸 강조한다.")

    s = site_slide(
        prs, nx(), TOTAL, "Part 3 · 기술 스택 ①", "vercel.com",
        "이 앱이 사는 집 — 배포(Deployment)",
        "“내가 만든 코드를, 전 세계 누구나 접속할 수 있는 인터넷 주소로 바꿔주는 곳”",
        [
            ("GitHub에 코드를 올리면, Vercel이 자동으로 감지해서 새로 빌드하고 배포한다", 0),
            ("그 결과가 recipe4fridge-pic.vercel.app 같은 실제 주소로 나온다", 0),
            ("무료 요금제로 개인 프로젝트를 시작하기에 충분하다", 0),
            ("우리는 \"배포\"라는 어려운 단어 대신, \"저장 버튼을 누르면 자동으로 세상에 공개된다\"고 "
             "기억하면 된다", 0, {"bold": True}),
        ],
        "Vercel을 \"내 코드가 실제로 사는 집\"이라고 소개한다. GitHub push 한 번에 자동 배포되는 흐름을 "
        "\"편지를 우체통에 넣으면 자동으로 상대방 집까지 배달되는 것\"처럼 설명해도 좋다. 이 강의를 준비하며 "
        "실제로 이 사이트를 통해 수십 차례 배포했다는 걸 이야기해준다.",
        CORAL,
    )

    s = site_slide(
        prs, nx(), TOTAL, "Part 3 · 기술 스택 ②", "supabase.com",
        "이 앱의 기억과 창고 — 회원·데이터베이스·사진 저장",
        "“나만의 작은 서버실을, 코드 한 줄 없이 클릭 몇 번으로 공짜로 빌리는 곳”",
        [
            ("회원가입/로그인(Auth) — 이메일·비밀번호, 구글 로그인을 대신 관리해준다", 0),
            ("데이터베이스(Postgres) — 사용자, 레시피, 평가 같은 표(table)들을 저장한다", 0),
            ("파일 저장소(Storage) — 업로드한 냉장고 사진을 안전하게 보관한다", 0),
            ("\"누가 어떤 데이터를 볼 수 있는가\"를 정하는 보안 규칙(RLS)도 여기서 설정한다", 0),
        ],
        "Supabase는 이 앱에서 가장 자주 등장하는 서비스다. 로그인, 데이터, 사진까지 세 가지 역할을 한 "
        "곳에서 다 한다는 것을 강조한다. \"예전 같으면 서버를 직접 빌리고 데이터베이스를 설치해야 했는데, "
        "지금은 가입만 하면 이 모든 게 준비된다\"는 것이 얼마나 큰 변화인지 짚어준다.",
        NAVY,
    )

    s = site_slide(
        prs, nx(), TOTAL, "Part 3 · 기술 스택 ③", "openrouter.ai",
        "이 앱의 AI 두뇌 렌탈 서비스",
        "“구글, 메타, 알리바바 같은 여러 회사의 AI 모델을 한 창구에서 골라 빌려 쓰는 곳”",
        [
            ("사진을 보고 재료를 알아내는 \"눈\"(비전 모델)을 빌려온다", 0),
            ("재료로 레시피를 써내는 \"글솜씨\"(텍스트 모델)를 빌려온다", 0),
            ("무료(:free) 모델도 제공해서, 비용 걱정 없이 배우고 실험할 수 있다", 0),
            ("한 모델이 응답을 못 하면 다른 모델로 자동으로 갈아타는 안전장치도 여기서 만든다", 0),
        ],
        "OpenRouter를 \"AI 인력 파견 업체\"에 비유해도 좋다. 여러 회사가 만든 AI 모델을 직접 하나하나 "
        "계약할 필요 없이, 한 창구에서 골라 쓸 수 있다는 편리함을 강조한다. 특히 무료 모델이 존재한다는 "
        "사실이 신입생들에게 \"나도 지금 당장 시작할 수 있다\"는 실질적 자신감을 준다.",
        CORAL_DARK,
    )

    s = site_slide(
        prs, nx(), TOTAL, "Part 3 · 기술 스택 ④", "console.cloud.google.com",
        "구글 로그인을 위한 신분증 발급소",
        "“우리 앱이 \\'구글 로그인\\'이라는 버튼을 달 수 있게, 구글이 발급해주는 열쇠(OAuth) 창구”",
        [
            ("여기서 \"이 앱은 이런 이름과 주소를 가진 진짜 프로젝트\"라고 등록한다", 0),
            ("등록하면 클라이언트 ID/키를 발급받아, 우리 앱에 \"구글로 로그인\" 버튼을 연결할 수 있다", 0),
            ("사용자 입장에서는 비밀번호를 새로 안 만들어도 되니 훨씬 편해진다", 0),
            ("네 사이트 중 가장 낯선 이름이지만, 역할은 \"신분증 발급\" 하나로 단순하다", 0, {"bold": True}),
        ],
        "Google Cloud Console은 이름이 가장 위압적이지만 우리가 쓰는 역할은 딱 하나, OAuth(구글 로그인) "
        "열쇠 발급이라는 것을 분명히 한다. \"이 큰 사이트 안에서 우리는 아주 작은 방 하나만 빌려 쓰는 것\"이라고 "
        "안심시켜준다.",
        GREEN,
    )

    s = content_slide(prs, "Part 3 · 기술 스택", "4개 사이트, 한 장으로 정리", nx(), TOTAL)
    table = [
        ("사이트", "한 줄 역할", "비유"),
        ("vercel.com", "코드를 실제 인터넷 주소로 배포", "공연장"),
        ("supabase.com", "로그인·데이터베이스·사진 저장소", "창고 겸 대기실"),
        ("openrouter.ai", "여러 회사의 AI 모델을 빌려 호출", "AI 섭외처"),
        ("console.cloud.google.com", "구글 로그인용 열쇠(OAuth) 발급", "신분증 발급소"),
    ]
    col_w = [Inches(3.6), Inches(5.4), Inches(2.8)]
    top = Inches(2.1)
    row_h = Inches(0.75)
    x0 = Inches(0.7)
    for r, row in enumerate(table):
        x = x0
        is_head = r == 0
        for c, text in enumerate(row):
            w = col_w[c]
            fill = CORAL if is_head else (WHITE if r % 2 else CREAM2)
            add_rect(s, x, top + row_h * r, w, row_h - Pt(2), fill=fill, line_color=LINE, line_w=0.75)
            add_text(s, text, x + Inches(0.15), top + row_h * r, w - Inches(0.3), row_h - Pt(2),
                      size=15 if not is_head else 16, color=WHITE if is_head else INK,
                      bold=is_head, anchor=MSO_ANCHOR.MIDDLE)
            x += w
    set_notes(s, "이 표를 사진 찍어두라고 권해도 좋을 만큼, 오늘 Part3의 핵심을 압축한 슬라이드다. "
                 "표 아래에 시간을 조금 남겨 질문을 받는다.")

    s = content_slide(prs, "Part 3 · 기술 스택", "실제로 가입하면 만나는 화면들", nx(), TOTAL)
    add_bullets(s, [
        ("vercel.com — GitHub 계정으로 로그인 → \"New Project\" → 저장소 선택 → 배포 버튼 클릭", 0),
        ("supabase.com — \"New Project\" → 프로젝트 이름/비밀번호 설정 → 몇 분 뒤 나만의 데이터베이스 완성", 0),
        ("openrouter.ai — 가입 후 \"Keys\" 메뉴에서 API 키 발급 → 무료 모델부터 바로 테스트 가능", 0),
        ("console.cloud.google.com — \"OAuth 동의 화면\" 만들기 → \"사용자 인증 정보\" → 클라이언트 ID 발급", 0),
        ("네 곳 모두 공통점: 회원가입 → 새 프로젝트 만들기 → 열쇠(키) 하나 발급받기, 이 패턴의 반복이다", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.4), size=18)
    set_notes(s, "가능하면 이 슬라이드에서 실제 화면을 잠깐 보여주거나(사전 녹화 화면 캡처), 다음 실습 "
                 "시간에 직접 가입해보는 과제로 이어준다. \"모든 사이트가 결국 같은 패턴\"이라는 걸 깨닫는 "
                 "순간이 신입생들에게 큰 안도감을 준다.")

    # ---------------- PART 4 ----------------
    section_divider(prs, 4, "개발 과정 실제로 따라가보기", "찰칵레시피가 오늘 모습이 되기까지", 35)
    nx()

    s = content_slide(prs, "Part 4 · 개발 과정", "이정표 — 계획했던 것과 실제로 만든 것", nx(), TOTAL)
    add_image_fit(s, os.path.join(DIAGRAMS, "milestones.png"), Inches(0.7), Inches(1.95),
                  Inches(11.9), Inches(4.8))
    set_notes(s, "처음 8단계 계획(파란 카드)은 전부 완료됐고, 그 뒤로 계획에 없던 확장 작업(보라 카드)이 "
                 "6개나 더 생겼다는 것을 짚어준다. 계획대로만 되는 프로젝트는 없으며, 오히려 대화를 "
                 "나누다 보면 처음엔 안 보이던 필요가 계속 발견된다는 것이 이번 파트의 핵심 메시지다.")

    s = process_slide(
        prs, nx(), TOTAL, "1", "PRD 먼저 쓰기 — 무엇을 만들지 AI와 함께 정리한다",
        "코드보다 먼저 \"문서\"부터 썼다. 목표, 화면 흐름, 데이터까지 말로 정리하고 나서야 코드를 시작했다.",
        [
            ("\"냉장고 사진으로 레시피를 추천하는 앱을 만들고 싶어\"라는 한 문장에서 출발", 0),
            ("클로드 코드와 대화하며 목표, 사용자 흐름, 기능 목록, 데이터 구조를 문서(PRD)로 정리", 0),
            ("이 문서가 이후 모든 대화의 \"공통 기준점\"이 되어준다", 0),
            ("설계도 없이 집을 짓지 않듯, 코드도 마찬가지다", 0, {"bold": True, "color": CORAL_DARK}),
        ],
        "PRD(제품 요구사항 문서)라는 용어를 처음 접하는 학생들에게는 \"만들기 전에 정리한 메모\" 정도로 "
        "쉽게 풀어준다. 이 문서가 이번 강의 자료(docs/PRD.md)로 실제 저장소에 남아있다는 것도 보여줄 수 "
        "있다.",
    )

    s = process_slide(
        prs, nx(), TOTAL, "2", "가짜(mock)로 먼저 전체 그림을 완성한다",
        "진짜 AI API를 연결하기 전에, 가짜 응답을 반환하는 \"모형\"으로 화면 전체 흐름부터 완성했다.",
        [
            ("이유: API 키가 아직 없어도, 화면과 흐름이 맞는지부터 검증할 수 있다", 0),
            ("사진 업로드 → (가짜)재료 인식 → (가짜)레시피 추천까지 끝까지 눌러볼 수 있는 \"뼈대\"를 먼저 완성", 0),
            ("뼈대가 튼튼하면, 나중에 가짜를 진짜로 바꾸는 일은 훨씬 쉬워진다", 0, {"bold": True}),
            ("이 지혜는 코딩을 넘어 모든 일에 통한다 — 완벽한 부품보다, 먼저 전체를 얼기설기 이어보기", 0),
        ],
        "이 단계가 신입생에게 특히 중요한 교훈이다. \"완벽하게 하나씩 끝내고 다음으로\"가 아니라 "
        "\"어설프더라도 처음부터 끝까지 먼저 이어보기\"라는 전략을, 이후 어떤 프로젝트에도 쓸 수 있는 "
        "지혜로 강조한다.",
        image=os.path.join(DIAGRAMS, "architecture_mock.png"),
    )

    s = process_slide(
        prs, nx(), TOTAL, "3", "진짜 AI로 교체한다",
        "가짜 응답을 지우고, openrouter.ai의 실제 비전·텍스트 모델을 연결했다.",
        [
            ("mock 코드는 지우지 않고 남겨뒀다 — \"가짜 로직 vs 실제 API\"를 비교하는 수업 자료로도 쓴다", 0),
            ("한 모델이 한도 초과나 오류로 응답을 못 하면, 다음 모델로 자동으로 갈아타는 안전장치를 추가", 0),
            ("실제 AI 응답은 형식이 들쭉날쭉할 수 있어, 응답을 관대하게 해석하는 처리도 더했다", 0),
            ("\"완벽한 AI\"를 기다리지 않고, 여러 겹의 안전망으로 현실적인 신뢰도를 만들었다", 0,
             {"bold": True, "color": CORAL_DARK}),
        ],
        "이 단계에서 \"AI도 가끔 틀리거나 응답을 못 할 수 있다\"는 사실을 솔직하게 알려준다. 그래서 "
        "\"대체 모델로 자동 전환\"이라는 안전장치를 만든 것이 왜 중요한지 실무적으로 설명한다.",
        image=os.path.join(DIAGRAMS, "architecture.png"),
    )

    s = process_slide(
        prs, nx(), TOTAL, "4", "배포하고, 실제로 써보고, 고친다",
        "Vercel에 배포한 뒤에도 끝이 아니다 — 실제 화면에서 눌러보며 계속 다듬었다.",
        [
            ("GitHub에 코드를 올리면 Vercel이 자동으로 새 버전을 배포한다", 0),
            ("배포 후 실제 사진을 올려 \"비전 API → 재료 인식 → 텍스트 API → 레시피 생성\"이 진짜로 "
             "도는지 직접 확인했다", 0),
            ("드롭다운 UI가 화면 밖으로 삐져나오는 것 같은 작은 문제도 실제로 써봐야 보인다", 0),
            ("\"다 만들었다\"가 아니라 \"써보니 이건 고쳐야겠다\"의 반복이 진짜 개발이다", 0, {"bold": True}),
        ],
        "이 단계에서는 실제로 이번 강의 자료를 준비하며 배포된 사이트에 직접 접속해 사진을 올려 "
        "테스트했던 경험을 생생하게 들려준다. \"이론이 아니라 실제로 만져봐야 문제가 보인다\"는 점을 "
        "강조한다.",
    )

    s = process_slide(
        prs, nx(), TOTAL, "5", "대화하다 보면, 계획에 없던 것도 생긴다",
        "PRD의 8단계 계획을 다 마친 뒤에도, 대화를 계속하며 6개의 새로운 작업이 자연스럽게 생겨났다.",
        [
            ("\"관리자만 보이는 이름을 바꿔야겠다\" → 브랜드 리뉴얼(찰칵레시피 이름과 로고)로 확장", 0),
            ("\"다른 회원 레시피도 볼 수 있으면 좋겠다\" → 관리자 전용 화면을 전체 공개 구조로 재설계", 0),
            ("\"AI 모델끼리 성능을 비교해볼 수 있을까?\" → AI 모델 품질 평가 시스템을 아예 새로 설계", 0),
            ("계획은 나침반이지, 감옥이 아니다 — 대화가 곧 기획의 연장이다", 0,
             {"bold": True, "color": CORAL_DARK}),
        ],
        "이 슬라이드는 \"바이브 코딩은 계획을 안 세운다\"는 오해를 바로잡는다. 계획(PRD)은 분명히 있었지만, "
        "실제 대화 과정에서 더 나은 아이디어가 계속 떠올랐고, 그걸 유연하게 반영할 수 있었던 것이 바이브 "
        "코딩의 강점이라는 걸 짚어준다.",
    )

    s = content_slide(prs, "Part 4 · 개발 과정", "미니 사례 — AI가 AI를 평가하게 만들다", nx(), TOTAL)
    add_image_fit(s, os.path.join(DIAGRAMS, "model_eval.png"), Inches(0.7), Inches(1.85),
                  Inches(11.9), Inches(4.9))
    set_notes(s, "이 사례는 오늘 강의에서 가장 \"고급스러운\" 예시다. 사용자가 1~5점으로 평가하는 것에 "
                 "더해, AI 모델 하나를 \"채점관\"으로 세워 다른 모델들의 결과를 자동으로 채점하게 만들었다. "
                 "이렇게 모인 점수가 실제로 화면의 모델 선택 순서를 자동으로 바꾼다는 것까지 설명하면, "
                 "\"AI를 활용해 AI를 관리한다\"는 다음 단계의 아이디어를 신입생들에게 심어줄 수 있다.")

    s = content_slide(prs, "Part 4 · 개발 과정", "데이터 모델도 결국 상식이다", nx(), TOTAL)
    add_image_fit(s, os.path.join(DIAGRAMS, "erd.png"), Inches(0.7), Inches(1.85),
                  Inches(11.9), Inches(4.9))
    set_notes(s, "복잡해 보이는 데이터베이스 설계도(ERD)도, 실은 우리가 일상에서 쓰는 \"표\"들의 관계일 "
                 "뿐이라는 걸 강조한다. 회원 한 명이 여러 세션을 갖고, 세션 하나가 여러 사진과 재료를 갖고, "
                 "그 결과로 레시피가 생기는 흐름을 화살표를 손가락으로 짚어가며 설명하면 좋다. 이것도 "
                 "결국 클로드 코드와의 대화로 설계됐다는 걸 상기시킨다.")

    # ---- BREAK 2 ----
    s = new_slide(prs, bg=NAVY)
    add_text(s, "마지막 쉬는 시간입니다", Inches(0.9), Inches(2.9), Inches(11), Inches(1.0), size=40,
              bold=True, color=WHITE)
    add_text(s, "10분 후, 오늘 강의의 결론 — \"여러분도 할 수 있다\"로 마무리합니다", Inches(0.9),
              Inches(4.0), Inches(11.3), Inches(0.6), size=18, color=CREAM2)
    nx()

    # ---------------- PART 5 ----------------
    section_divider(prs, 5, "여러분도 할 수 있다는 확신", "\"나는 코드를 몰라서\"는 이제 핑계가 안 된다", 20)
    nx()

    s = content_slide(prs, "Part 5 · 자신감", "\"나는 코드를 몰라서 못한다\"는 착각", nx(), TOTAL)
    add_bullets(s, [
        ("오늘 3시간 동안 보여드린 화면 흐름, 아키텍처, 데이터 모델 — 전부 문장으로 설명해서 만들었다", 0),
        ("문법을 몰라도, \"무엇을 원하는지\"를 구체적으로 말할 수 있으면 시작할 수 있다", 0),
        ("처음엔 완벽한 문장이 아니어도 된다 — 대화를 주고받으며 점점 정확해지면 된다", 0),
        ("진짜 필요한 것은 문법 지식이 아니라, 좋은 질문을 던지는 습관이다", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "이 슬라이드에서 신입생들의 가장 큰 심리적 장벽(\"나는 전공자가 아니라서\")을 정면으로 "
                 "다룬다. 오늘 다룬 모든 예시 문장이 실제로 평범한 한국어 문장이었다는 걸 다시 짚어준다.")

    s = content_slide(prs, "Part 5 · 자신감", "오늘 앱도 결국 대화의 누적이었다", nx(), TOTAL)
    add_bullets(s, [
        ("\"냉장고 사진으로 레시피 추천하는 앱 만들고 싶어\" — 씨앗이 된 첫 문장", 0),
        ("\"사진이 더 크면 좋겠어\", \"확대 기능을 넣어도 좋고\" — 써보고 나서 나온 개선 요청", 0),
        ("\"이 판정도 무료 모델을 쓰면 어떨까?\" — 비용을 고민하다 나온 아이디어", 0),
        ("작은 문장 수백 개가 쌓여서, 지금의 찰칵레시피가 되었다", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "실제로 이 강의 자료를 준비하며 오갔던 대화 문장들을 그대로 인용한 것임을 밝히면 "
                 "설득력이 커진다. \"큰 그림을 한 번에 그린 게 아니라, 작은 대화가 계속 쌓인 것\"이라는 "
                 "메시지를 강조한다.")

    s = content_slide(prs, "Part 5 · 자신감", "여러분이라면 무엇을 만들겠습니까?", nx(), TOTAL)
    add_bullets(s, [
        ("동아리 출석과 회비를 정리해주는 앱", 0),
        ("자취생을 위한 \"이번 주 장보기 리스트 추천\" 앱", 0),
        ("전공 수업 조모임 일정을 자동으로 맞춰주는 앱", 0),
        ("내 강아지/고양이 사진을 보고 오늘 컨디션을 기록해주는 앱", 0),
        ("이 목록의 공통점: 전부 \"내 주변의 작은 불편\"에서 출발한다는 것", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "이 슬라이드는 참여형으로 진행한다. 예시를 던진 뒤 \"여러분 주변엔 어떤 불편이 있나요?\" "
                 "라고 실제로 몇 명에게 질문해 짧게 답을 들어본다. 다음 시간 과제(자기 아이디어로 PRD "
                 "한 페이지 써오기) 예고로 이어가도 좋다.")

    s = content_slide(prs, "Part 5 · 자신감", "실패해도 괜찮은 이유", nx(), TOTAL)
    add_bullets(s, [
        ("모든 변경 이력은 기록되어 남는다(버전 관리, Git) — 문제가 생기면 이전 상태로 되돌릴 수 있다", 0),
        ("오늘 배운 것처럼 mock(가짜)으로 먼저 시도해보고, 안전하게 진짜로 바꿔갈 수 있다", 0),
        ("AI에게 \"이거 왜 이렇게 됐어?\"라고 물어보면, 원인을 함께 찾아준다", 0),
        ("\"완벽한 한 번\"이 아니라 \"작게, 자주 시도하고 되돌릴 수 있는 여러 번\"이 이 시대의 안전망이다", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "실패에 대한 두려움을 낮추는 슬라이드. 버전 관리(Git)라는 단어를 처음 듣는 학생들에게는 "
                 "\"모든 순간의 저장 파일이 자동으로 남아있어 언제든 되감기할 수 있는 게임의 세이브 포인트\"에 "
                 "비유해도 좋다.")

    # ---------------- PART 6 ----------------
    section_divider(prs, 6, "AI 활용 개발자에게 필요한 7가지 역량", "코딩 실력보다 오래가는 것", 25)
    nx()

    s = content_slide(prs, "Part 6 · 마무리", "왜 \"코딩 실력\"만으로는 부족해지는가", nx(), TOTAL)
    add_bullets(s, [
        ("문법을 외우고 빠르게 타이핑하는 능력은, AI가 이미 사람보다 잘한다", 0),
        ("그렇다면 이제 사람에게 남는 것, 오히려 더 중요해지는 것은 무엇인가", 0),
        ("오늘 찰칵레시피를 만드는 과정에서 실제로 계속 쓰였던 능력 7가지를 정리했다", 0,
         {"bold": True, "color": CORAL_DARK}),
        ("이 7가지는 전공과 무관하게, 지금 이 자리의 모든 신입생에게 이미 씨앗이 있는 능력들이다", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "Part6 전체의 문을 여는 슬라이드. \"코딩을 잘하는 사람\"에서 \"AI와 함께 좋은 것을 "
                 "만드는 사람\"으로 기준이 이동했다는 것을 분명히 선언한다.")

    competencies = [
        ("사람들의 욕구를 파악해서\n꿈을 키워내는 능력",
         "기술이 아니라 사람에서 출발한다",
         [
             "\"냉장고 사진 → 오늘 뭐 해먹지\"는 기술이 아니라 사람의 불편에서 시작됐다",
             "좋은 개발자는 코드보다 먼저, 사람이 진짜 원하는 것을 알아채는 사람이다",
             "작은 불편을 발견하는 눈이, 곧 좋은 아이디어의 씨앗이다",
         ]),
        ("성찰을 통한\n위치 파악 및 목표 설정 능력",
         "지금 어디에 있고, 어디로 가는지 알아야 한다",
         [
             "PRD를 먼저 쓴 것도 결국 \"우리가 지금 뭘 하려는지\"를 스스로 확인하는 과정이었다",
             "이정표(마일스톤)를 계속 되돌아보며 \"계획대로인가, 방향을 바꿔야 하나\"를 점검했다",
             "AI는 방향을 대신 정해주지 않는다 — 방향은 내가 정하고, AI는 그 길을 함께 걷는다",
         ]),
        ("상식적 판단력",
         "\"이게 말이 되나?\"를 끝까지 놓지 않는 힘",
         [
             "AI가 내놓은 결과를 \"그럴듯해 보인다\"로 끝내지 않고 \"정말 맞나?\"로 검증했다",
             "평가 점수가 5건도 안 쌓였는데 순위를 확정 짓지 않은 것도 상식적 판단이다",
             "AI의 속도에 상식이라는 브레이크를 다는 사람이, 신뢰할 수 있는 결과를 만든다",
         ]),
        ("클로드 코드 협업에 대한 신뢰",
         "동료를 믿고 일을 맡길 줄 아는 능력",
         [
             "모든 걸 직접 타이핑하려 하지 않고, 방향을 정한 뒤 실행은 믿고 맡겼다",
             "결과를 함께 검토하고, 아니다 싶으면 다시 요청하는 \"건강한 협업\"을 반복했다",
             "신뢰는 맹신이 아니다 — 확인하고, 고치고, 다시 신뢰하는 과정의 반복이다",
         ]),
        ("개선에 대한\n열망과 열정",
         "\"이 정도면 됐다\"에서 한 걸음 더",
         [
             "사진이 작아 보이면 확대 기능을, 이름이 밋밋하면 브랜드를 새로 고민했다",
             "완성됐다고 끝내지 않고, 써보면서 계속 \"더 나아질 부분\"을 찾아냈다",
             "이 열정이 있어야, 계획에 없던 6개의 확장 작업도 마다하지 않고 해낼 수 있었다",
         ]),
        ("포기하지 않는\n의지와 체력",
         "안 되면 될 때까지, 다르게 다시 시도하는 힘",
         [
             "무료 AI 모델이 오류를 내면, 될 때까지 다른 모델로 자동 전환하는 안전망을 만들었다",
             "화면이 깨지고, 배포가 실패하고, 파일 권한이 꼬여도 원인을 찾아 끝까지 고쳤다",
             "긴 작업일수록 \"오늘은 여기까지\"를 알고, 다시 이어갈 체력을 관리하는 것도 실력이다",
         ]),
        ("차별적 성과 창출",
         "남들과 같은 것을 만들지 않는 용기",
         [
             "인기 투표와 AI 모델 품질 평가를 분리한 설계처럼, 남이 안 하는 디테일까지 챙겼다",
             "\"AI가 AI를 평가하게 만드는\" 시스템처럼, 한 걸음 더 나아간 아이디어를 실제로 구현했다",
             "결국 도구는 누구에게나 열려 있다 — 차이를 만드는 건 그 도구로 무엇을 어디까지 밀어붙이는가다",
         ]),
    ]
    for idx, (title, essence, bullets) in enumerate(competencies, start=1):
        s = competency_slide(prs, nx(), TOTAL, idx, title, essence, bullets,
                              f"역량 {idx} 상세 설명. 반드시 오늘 강의에서 실제로 등장했던 구체적 장면과 "
                              f"연결지어 이야기한다. 추상적 훈화로 끝나지 않도록, \"오늘 우리가 이걸 이렇게 "
                              f"써먹었죠\"라는 방식으로 되짚어준다.")

    s = content_slide(prs, "Part 6 · 마무리", "7가지 역량, 한 장으로 정리", nx(), TOTAL)
    labels = [
        "① 욕구 파악 · 꿈을 키우는 능력", "② 성찰을 통한 위치 파악 · 목표 설정",
        "③ 상식적 판단력", "④ 클로드 코드 협업에 대한 신뢰",
        "⑤ 개선에 대한 열망 · 열정", "⑥ 포기하지 않는 의지 · 체력",
        "⑦ 차별적 성과 창출",
    ]
    cols = 2
    cw, ch = Inches(5.75), Inches(0.85)
    gx, gy = Inches(0.3), Inches(0.2)
    for i, label in enumerate(labels):
        r, c = divmod(i, cols)
        x = Inches(0.7) + c * (cw + gx)
        y = Inches(2.0) + r * (ch + gy)
        add_rect(s, x, y, cw, ch, fill=CREAM2, radius=0.15)
        add_text(s, label, x + Inches(0.25), y, cw - Inches(0.5), ch, size=17, color=INK, bold=True,
                  anchor=MSO_ANCHOR.MIDDLE)
    set_notes(s, "전체 요약 슬라이드. 이 슬라이드를 사진 찍어가라고 권해도 좋다. 7가지 중 오늘 강의에서 "
                 "가장 인상 깊었던 것을 한 가지씩 골라보라고 질문을 던지며 마무리 토론으로 이어갈 수 있다.")

    s = content_slide(prs, "Part 6 · 마무리", "오늘의 작은 과제", nx(), TOTAL)
    add_bullets(s, [
        ("① 내 주변의 작은 불편 하나를 문장으로 적어보기 (\"OOO 하기가 매번 귀찮다\")", 0),
        ("② 그 불편을 풀어주는 앱을 한 문단으로 설명해보기 — 이것이 나만의 첫 PRD가 된다", 0),
        ("③ vercel.com, supabase.com 중 한 곳에 가입만이라도 해보기", 0),
        ("완벽하지 않아도 된다. 오늘 배운 4단계 루프처럼, 다음 시간에 함께 다듬어가면 된다", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "다음 수업과의 연결고리. 과제를 거창하게 만들지 말고, 오늘 배운 태도(작게 시작해서 "
                 "대화로 키워가기)를 그대로 적용해보게 한다.")

    s = new_slide(prs, bg=CREAM)
    add_rect(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill=CORAL)
    add_text(s, "마치며", Inches(0.9), Inches(1.5), Inches(10), Inches(0.5), size=18, color=CORAL_DARK,
              bold=True)
    add_text(s, "무엇이든 만들 수 있습니다,\n여러분이 정말 원한다면", Inches(0.9), Inches(2.1),
              Inches(11.5), Inches(1.9), size=40, color=INK, bold=True, line_spacing=1.25)
    add_text(s,
             "저는 40년을 코드와 함께 살았지만, 오늘 여러분과 함께 본 이 3시간의 여정이야말로 지금 이 "
             "시대에 코딩을 배운다는 것의 진짜 의미라고 생각합니다. 냉장고 사진 한 장에서 시작한 작은 "
             "아이디어가, 실제로 접속할 수 있는 서비스가 되기까지, 문법이 아니라 대화가 그 다리를 놓았습니다.\n\n"
             "여러분에게도 이미 그 씨앗이 있습니다. 오늘 배운 것을 가지고, 여러분만의 첫 문장을 "
             "시작해보시기 바랍니다.",
             Inches(0.9), Inches(4.15), Inches(11.5), Inches(2.6), size=18, color=INK, line_spacing=1.4)
    set_notes(s, "마지막 인사. 진심을 담아 천천히 읽어준다. 박수를 유도하거나, 질의응답 시간을 이어서 "
                 "가질지는 현장 분위기에 맞춰 결정한다.")

    return prs, TOTAL, n


if __name__ == "__main__":
    prs, total, n = build()
    print(f"planned TOTAL={total}, n counter={n}, actual slide count={len(prs.slides)}")
    out_path = os.path.join(PROJECT_ROOT, "docs", "찰칵레시피_바이브코딩_3시간강의.pptx")
    prs.save(out_path)
    print("wrote", out_path)



