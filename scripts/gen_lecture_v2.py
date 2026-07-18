# -*- coding: utf-8 -*-
"""
3-hour+ lecture deck v2.0: "찰칵레시피 개발 여정 — 클로드 코드로 배우는 바이브 코딩"
Full technical deep-dive edition. Built with python-pptx, no network/template dependency.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

PROJECT_ROOT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..")
DIAGRAMS = os.path.join(PROJECT_ROOT, "docs", "diagrams")
CAPTURED = os.path.join(PROJECT_ROOT, "captured")
CODE_FONT = "Menlo"

FONT = "맑은 고딕"

CREAM = RGBColor(0xFB, 0xF4, 0xEC)
CREAM2 = RGBColor(0xF3, 0xE9, 0xDA)
INK = RGBColor(0x2C, 0x23, 0x1D)
MUTED = RGBColor(0x8A, 0x7B, 0x6E)
CORAL = RGBColor(0xD9, 0x64, 0x3A)
CORAL_DARK = RGBColor(0xB1, 0x4B, 0x26)
LINE = RGBColor(0xE8, 0xDC, 0xCB)
GREEN = RGBColor(0x3E, 0x81, 0x56)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x2E, 0x3A, 0x4E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PAGE = {"n": 0}


def new_slide(prs, bg=CREAM):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()
    bg_shape.shadow.inherit = False
    # send to back
    sp = bg_shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    PAGE["n"] += 1
    return slide


def _set_text(tf, text, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT,
              font=FONT, line_spacing=1.15, italic=False):
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color


def add_text(slide, text, left, top, width, height, size=20, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.15,
             italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    _set_text(tf, text, size=size, color=color, bold=bold, align=align, font=font,
              line_spacing=line_spacing, italic=italic)
    return box


def add_bullets(slide, items, left, top, width, height, size=18, color=INK, font=FONT,
                 space_after=10, line_spacing=1.2):
    """items: list of (text, level, opts) where opts is dict with optional
    'bold','color','bullet' keys. level 0/1/2 controls indent."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            text, level, opts = item, 0, {}
        elif len(item) == 2:
            text, level = item
            opts = {}
        else:
            text, level, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        indent = level * 0.35
        bullet_char = opts.get("bullet", ["●", "–", "·"][min(level, 2)])
        prefix = f"{bullet_char}  " if bullet_char else ""
        r = p.add_run()
        r.text = ("    " * level) + prefix + text
        r.font.size = Pt(opts.get("size", size - level * 1.5))
        r.font.bold = opts.get("bold", level == 0)
        r.font.name = font
        r.font.color.rgb = opts.get("color", color if level == 0 else MUTED)
    return box


def add_rect(slide, left, top, width, height, fill=None, line_color=None, line_w=1.0,
             radius=None, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if radius:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = shadow
    return shp


def add_image_fit(slide, path, left, top, max_w, max_h, align="center"):
    im = Image.open(path)
    ratio = im.width / im.height
    w = max_w
    h = w / ratio
    if h > max_h:
        h = max_h
        w = h * ratio
    if align == "center":
        left = left + (max_w - w) / 2
    pic = slide.shapes.add_picture(path, left, top + (max_h - h) / 2, width=int(w), height=int(h))
    return pic


def footer(slide, label, page_no, total):
    add_text(slide, label, Inches(0.6), Inches(7.08), Inches(8), Inches(0.35),
              size=10, color=MUTED, font=FONT)
    add_text(slide, f"{page_no} / {total}", Inches(12.2), Inches(7.08), Inches(0.9), Inches(0.35),
              size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT)


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def section_divider(prs, no, title, subtitle):
    slide = new_slide(prs, bg=CORAL_DARK)
    add_text(slide, f"PART {no}", Inches(0.9), Inches(2.5), Inches(4), Inches(0.6),
             size=20, color=CREAM2, bold=True, font=FONT)
    add_text(slide, title, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.6),
             size=40, color=WHITE, bold=True, font=FONT)
    add_text(slide, subtitle, Inches(0.9), Inches(4.3), Inches(11), Inches(1.0),
             size=18, color=CREAM2, font=FONT)
    return slide


def content_slide(prs, kicker, title, page_no, total, bg=CREAM):
    slide = new_slide(prs, bg=bg)
    if kicker:
        add_text(slide, kicker, Inches(0.7), Inches(0.42), Inches(10), Inches(0.4),
                  size=14, color=CORAL_DARK, bold=True, font=FONT)
    add_text(slide, title, Inches(0.7), Inches(0.75), Inches(11.9), Inches(1.0),
              size=30, color=INK, bold=True, font=FONT)
    add_rect(slide, Inches(0.7), Inches(1.62), Inches(1.1), Pt(3), fill=CORAL)
    footer(slide, "찰칵레시피로 배우는 바이브 코딩", page_no, total)
    return slide


def site_slide(prs, page_no, total, tag, domain, role_kr, analogy, bullets, notes, accent):
    slide = content_slide(prs, tag, domain, page_no, total)
    add_rect(slide, Inches(0.7), Inches(1.95), Inches(4.5), Inches(0.62), fill=accent, radius=0.5)
    add_text(slide, role_kr, Inches(0.7), Inches(1.95), Inches(4.5), Inches(0.62),
              size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, analogy, Inches(0.7), Inches(2.75), Inches(11.6), Inches(0.9),
              size=17, color=MUTED, italic=True)
    add_bullets(slide, bullets, Inches(0.7), Inches(3.7), Inches(11.8), Inches(3.1), size=18)
    set_notes(slide, notes)
    return slide


def competency_slide(prs, page_no, total, idx, title_kr, essence, bullets, notes):
    # A single consistent background across all 7 competency slides (per feedback:
    # alternating colors slide-to-slide read as noisy, not intentional).
    slide = new_slide(prs, bg=NAVY)
    add_text(slide, f"{idx:02d}", Inches(0.7), Inches(0.5), Inches(3), Inches(1.6),
              size=64, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    n_title_lines = title_kr.count("\n") + 1
    title_h = 0.65 * n_title_lines + 0.35
    add_text(slide, title_kr, Inches(0.7), Inches(2.0), Inches(11.8), Inches(title_h),
              size=32, color=WHITE, bold=True, line_spacing=1.15)
    essence_top = 2.0 + title_h + 0.2
    add_text(slide, essence, Inches(0.7), Inches(essence_top), Inches(11.6), Inches(0.6),
              size=18, color=CREAM2, italic=True)
    bullets_top = essence_top + 0.75
    add_bullets(slide, [(t, 0, {"color": CREAM2, "bold": False}) for t in bullets],
                Inches(0.7), Inches(bullets_top), Inches(11.8), Inches(7.0 - bullets_top), size=18,
                color=CREAM2)
    footer(slide, "AI 활용 개발자에게 필요한 7가지 역량", page_no, total)
    set_notes(slide, notes)
    return slide


def process_slide(prs, page_no, total, step_no, title, desc, bullets, notes, image=None):
    slide = content_slide(prs, f"개발 과정 {step_no}단계", title, page_no, total)
    body_w = Inches(6.6) if image else Inches(11.8)
    add_text(slide, desc, Inches(0.7), Inches(1.95), body_w, Inches(0.9), size=17, color=MUTED, italic=True)
    add_bullets(slide, bullets, Inches(0.7), Inches(2.85), body_w, Inches(3.9), size=17)
    if image:
        add_image_fit(slide, image, Inches(7.5), Inches(1.95), Inches(5.15), Inches(4.6))
    footer(slide, "찰칵레시피로 배우는 바이브 코딩", page_no, total)
    set_notes(slide, notes)
    return slide


def toc_slide(prs, page_no, total, title, parts, footer_label):
    """Table of contents / roadmap slide — no time estimates (per feedback)."""
    slide = content_slide(prs, "오늘의 순서", title, page_no, total)
    top = Inches(2.0)
    row_h = Inches(0.62)
    for i, (part_no, name) in enumerate(parts):
        y = top + row_h * i
        bg = CREAM2 if i % 2 == 0 else RGBColor(0xFA, 0xF3, 0xE8)
        add_rect(slide, Inches(0.7), y, Inches(11.8), row_h - Pt(4), fill=bg, radius=0.15)
        add_text(slide, part_no, Inches(0.95), y, Inches(1.6), row_h - Pt(4), size=15,
                  color=CORAL_DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, name, Inches(2.5), y, Inches(9.7), row_h - Pt(4), size=16, color=INK,
                  anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, footer_label, page_no, total)
    return slide


def site_step_slide(prs, page_no, total, tag_label, site_title, subtitle, steps, result_box, notes,
                     accent=CORAL):
    """A numbered, follow-along walkthrough for one concrete task on one site.
    steps: list of (step_title, step_detail) strings.
    result_box: (heading, text) shown as a highlighted "이제 이걸 갖게 됩니다" callout, or None.
    """
    slide = content_slide(prs, tag_label, site_title, page_no, total)
    if subtitle:
        add_text(slide, subtitle, Inches(0.7), Inches(1.9), Inches(11.6), Inches(0.5), size=15,
                  color=MUTED, italic=True)
    top = Inches(2.5) if subtitle else Inches(2.0)
    avail_bottom = Inches(6.85)
    n = len(steps)
    result_h = Inches(0.85) if result_box else Inches(0)
    row_area = avail_bottom - top - result_h - (Inches(0.15) if result_box else Inches(0))
    row_h = Emu(int(row_area / n)) if n else row_area
    for i, (stitle, sdetail) in enumerate(steps):
        y = top + row_h * i
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Pt(2), Inches(0.34), Inches(0.34))
        circ.fill.solid()
        circ.fill.fore_color.rgb = accent
        circ.line.fill.background()
        circ.shadow.inherit = False
        tf = circ.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT
        add_text(slide, stitle, Inches(1.25), y, Inches(3.6), row_h, size=16, color=INK, bold=True,
                  anchor=MSO_ANCHOR.TOP)
        add_text(slide, sdetail, Inches(4.95), y, Inches(7.55), row_h, size=14.5, color=MUTED,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.25)
    if result_box:
        ry = avail_bottom - result_h
        add_rect(slide, Inches(0.7), ry, Inches(11.8), result_h, fill=CREAM2, radius=0.12)
        heading, text = result_box
        add_text(slide, f"✅ {heading}", Inches(0.95), ry + Pt(6), Inches(11.3), Inches(0.3), size=13.5,
                  color=CORAL_DARK, bold=True)
        add_text(slide, text, Inches(0.95), ry + Inches(0.36), Inches(11.3), result_h - Inches(0.4),
                  size=13.5, color=INK, line_spacing=1.2)
    set_notes(slide, notes)
    return slide


def code_slide(prs, page_no, total, tag_label, title, file_path, code_lines, caption, notes,
                bullets=None):
    """Show a real, short code excerpt with its file path/GitHub-style location."""
    slide = content_slide(prs, tag_label, title, page_no, total)
    add_text(slide, f"📄 {file_path}", Inches(0.7), Inches(1.9), Inches(11.6), Inches(0.4), size=14,
              color=CORAL_DARK, bold=True, font=CODE_FONT)
    code_w = Inches(7.5) if bullets else Inches(11.8)
    box = add_rect(slide, Inches(0.7), Inches(2.35), code_w, Inches(3.55), fill=INK, radius=0.05)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.25)
    tf.margin_top = tf.margin_bottom = Inches(0.2)
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = line
        r.font.size = Pt(13)
        r.font.name = CODE_FONT
        r.font.color.rgb = RGBColor(0xE8, 0xE3, 0xD8)
    if bullets:
        add_bullets(slide, bullets, Inches(8.4), Inches(2.35), Inches(4.15), Inches(3.55), size=14.5)
    add_text(slide, caption, Inches(0.7), Inches(6.05), Inches(11.8), Inches(0.75), size=13.5,
              color=MUTED, italic=True, line_spacing=1.25)
    footer(slide, "찰칵레시피로 배우는 바이브 코딩", page_no, total)
    set_notes(slide, notes)
    return slide


def screenshot_box(slide, left, top, width, height, filename, placeholder_hint):
    """Embed captured/<filename> if it exists; otherwise draw a dashed placeholder
    telling the presenter exactly which live screen to capture and drop in later."""
    path = os.path.join(CAPTURED, filename)
    if os.path.exists(path):
        add_rect(slide, left, top, width, height, fill=WHITE, line_color=LINE, line_w=1.0, radius=0.04)
        add_image_fit(slide, path, left + Inches(0.06), top + Inches(0.06),
                      width - Inches(0.12), height - Inches(0.12))
    else:
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        rect.adjustments[0] = 0.04
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xF4, 0xEF, 0xE6)
        rect.line.color.rgb = MUTED
        rect.line.width = Pt(1.25)
        rect.line.dash_style = None
        add_text(slide, "📸 스크린샷 자리", left + Inches(0.25), top + Inches(0.2), width - Inches(0.5),
                  Inches(0.4), size=14, color=MUTED, bold=True)
        add_text(slide, placeholder_hint, left + Inches(0.25), top + Inches(0.65), width - Inches(0.5),
                  height - Inches(0.9), size=12.5, color=MUTED, line_spacing=1.3)
        add_text(slide, f"captured/{filename}", left + Inches(0.25), top + height - Inches(0.4),
                  width - Inches(0.5), Inches(0.32), size=10.5, color=MUTED, italic=True)


def case_study_slide(prs, page_no, total, case_no, user_quote, diagnosis, fix, lesson, lesson_kind,
                      shot_filename, shot_hint, notes):
    """요청 -> 진단 -> 해결 -> 스크린샷 -> 배운 점, 카드형 사례 슬라이드.
    lesson_kind: '기술' | '디자인' | '제품' | '방법론' 라벨로 표시."""
    slide = content_slide(prs, f"사례 {case_no}", "사용자 테스트에서 나온 실제 요청", page_no, total)
    add_rect(slide, Inches(0.7), Inches(1.9), Inches(7.3), Inches(0.95), fill=CREAM2, radius=0.1)
    add_text(slide, f"“{user_quote}”", Inches(0.95), Inches(1.9), Inches(6.8), Inches(0.95), size=14.5,
              color=INK, italic=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    add_text(slide, "진단", Inches(0.7), Inches(2.98), Inches(1.2), Inches(0.32), size=13, color=CORAL_DARK,
              bold=True)
    add_text(slide, diagnosis, Inches(0.7), Inches(3.3), Inches(7.3), Inches(1.05), size=13, color=INK,
              line_spacing=1.2)
    add_text(slide, "해결", Inches(0.7), Inches(4.42), Inches(1.2), Inches(0.32), size=13, color=CORAL_DARK,
              bold=True)
    add_text(slide, fix, Inches(0.7), Inches(4.74), Inches(7.3), Inches(1.2), size=13, color=INK,
              line_spacing=1.2)
    add_rect(slide, Inches(0.7), Inches(6.0), Inches(7.3), Inches(1.0), fill=NAVY, radius=0.1)
    add_text(slide, f"배운 점 ({lesson_kind})", Inches(0.95), Inches(6.08), Inches(6.9), Inches(0.28),
              size=11.5, color=CREAM2, bold=True)
    add_text(slide, lesson, Inches(0.95), Inches(6.38), Inches(6.9), Inches(0.58), size=11.5, color=WHITE,
              line_spacing=1.1)
    screenshot_box(slide, Inches(8.25), Inches(1.9), Inches(4.25), Inches(5.1), shot_filename, shot_hint)
    footer(slide, "실사용 테스트와 대응 사례", page_no, total)
    set_notes(slide, notes)
    return slide


def worksheet_slide(prs, page_no, total, title, prompts, notes):
    """A fill-in-the-blank style worksheet the student can literally answer during
    the lecture (or as homework) to draft their own PRD."""
    slide = content_slide(prs, "실습 워크시트", title, page_no, total)
    top = Inches(2.0)
    row_h = Inches(0.95)
    for i, (q, hint) in enumerate(prompts):
        y = top + row_h * i
        add_rect(slide, Inches(0.7), y, Inches(11.8), row_h - Pt(6), fill=WHITE, line_color=LINE,
                  line_w=1.0, radius=0.08)
        add_text(slide, f"Q{i + 1}. {q}", Inches(0.95), y + Pt(6), Inches(11.3), Inches(0.35), size=15,
                  color=INK, bold=True)
        add_text(slide, hint, Inches(0.95), y + Inches(0.42), Inches(11.3), Inches(0.45), size=12.5,
                  color=MUTED, italic=True)
    footer(slide, "PRD 작성 과정", page_no, total)
    set_notes(slide, notes)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    TOTAL = 101

    n = 0

    def nx():
        nonlocal n
        n += 1
        return n

    # ================= 0. 표지 & 전체 여정 =================
    s = new_slide(prs, bg=CREAM)
    add_rect(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill=CORAL)
    add_text(s, "대학 신입생을 위한 심화 실습형 강의 · v2.0", Inches(0.9), Inches(1.35), Inches(11), Inches(0.5),
              size=17, color=CORAL_DARK, bold=True)
    add_text(s, "찰칵레시피 개발 여정, 완전판", Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.3),
              size=44, color=INK, bold=True)
    add_text(s, "클로드 코드로 배우는 바이브 코딩\n— 구상, 계정 준비, 개발, 배포, 테스트, 보완까지 전부",
              Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.3), size=22, color=INK, line_spacing=1.3)
    add_text(s, "냉장고 사진 한 장 → 오늘 뭐 해먹지, 그 답을 만든 이야기 (v1 3시간 압축판의 확장·심화판)",
              Inches(0.9), Inches(4.3), Inches(11), Inches(0.6), size=15, color=MUTED, italic=True)
    add_rect(s, Inches(0.9), Inches(5.9), Inches(11.5), Pt(1.2), fill=LINE)
    add_text(s, "40년차 개발자가 신입생에게 건네는, 처음부터 끝까지 따라 하는 코딩 수업", Inches(0.9),
              Inches(6.1), Inches(11.5), Inches(0.5), size=14, color=MUTED)
    n = 1

    s = content_slide(prs, "시작하기 전에", "전체 개발 여정 지도 — 오늘 우리가 걸을 열 걸음", nx(), TOTAL)
    add_image_fit(s, os.path.join(DIAGRAMS, "journey.png"), Inches(0.7), Inches(1.9), Inches(11.9),
                  Inches(4.9))
    set_notes(s, "오늘 강의의 지도가 되는 슬라이드다. 아이디어에서 시작해 PRD를 쓰고, 5개 사이트에 계정을 "
                 "준비하고, mock으로 전체 흐름을 먼저 완성한 뒤 실제 AI로 교체하고, 배포하고, 실사용 테스트를 "
                 "거쳐 보완하고, 계획에 없던 확장을 하고, 마지막에 문서화까지 — 이 10단계를 오늘 각 파트에서 "
                 "하나씩 훨씬 상세하게 다룰 것이라고 예고한다. 특히 4~8단계는 한 번에 끝나지 않고 여러 번 "
                 "되돌아간다는 점을 강조해, 학생들이 실패를 반복으로 자연스럽게 받아들이게 한다.")

    s = toc_slide(prs, nx(), TOTAL, "오늘의 순서", [
        ("Part 1", "우리가 만든 것 — 찰칵레시피"),
        ("Part 2", "바이브 코딩이란 무엇인가"),
        ("Part 3", "PRD 작성 과정 상세 안내"),
        ("Part 4", "개발 착수 전 준비 — 5개 사이트 실습 가이드"),
        ("Part 5", "개발 과정 기술 심화 — mock에서 실전까지"),
        ("Part 6", "테스트하고 보완한 이야기 — 실제 사용자 피드백과 대응"),
        ("Part 7", "여러분도 할 수 있다는 확신"),
        ("Part 8", "AI 활용 개발자에게 필요한 7가지 역량"),
    ], "찰칵레시피로 배우는 바이브 코딩 · 완전판")
    set_notes(s, "이번 v2.0은 시간 배분표를 따로 두지 않는다. 내용이 훨씬 방대해져서 한 번의 3시간 세션보다는 "
                 "여러 회차 워크숍이나 실습 수업으로 나눠 진행하는 편이 더 어울리기 때문이다. 강의자가 현장 "
                 "상황에 맞게 파트별로 시간을 유연하게 배분하면 된다.")

    s = content_slide(prs, "시작하기 전에", "이 강의를 들으며 그대로 따라 하려면", nx(), TOTAL)
    add_bullets(s, [
        ("노트북과 인터넷 연결을 준비합니다.", 0),
        ("GitHub, Vercel, Supabase, OpenRouter, Google Cloud Console 계정을 미리 만들어 두면 실습이 "
         "빨라집니다. 강의 중에 만들어도 늦지 않습니다.", 0),
        ("강의 중 나오는 화면, 코드, 다이어그램은 전부 실제로 배포된 찰칵레시피 저장소 그대로입니다. 지어낸 "
         "예시가 아닙니다.", 0),
        ("완벽히 이해하지 못해도 괜찮습니다. 오늘의 목표는 세부 문법 암기가 아니라 \"전체 그림\"을 보고, "
         "여러분의 프로젝트에 그대로 옮겨 쓸 수 있는 순서를 익히는 것입니다.", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.4), size=18)
    set_notes(s, "실습형 강의임을 분명히 하고, 완벽주의에 대한 부담을 낮춰준다. 이 슬라이드에서 미리 계정 "
                 "생성을 독려하되, 못했어도 진도를 막지 않는다고 안심시킨다.")

    # ================= PART 1 =================
    section_divider(prs, 1, "우리가 만든 것", "찰칵레시피 — 냉장고 사진 한 장의 힘")
    nx()

    s = content_slide(prs, "Part 1 · 우리가 만든 것", "찰칵레시피는 어떤 문제를 풀어주는 앱인가?", nx(), TOTAL)
    add_bullets(s, [
        ("모두가 겪는 아주 사소하지만 매일 반복되는 고민이 있습니다. 바로 \"오늘 뭐 해먹지?\"입니다.", 0),
        ("냉장고 문을 열어보면 재료는 있는데, 그걸로 뭘 만들 수 있는지는 막상 떠오르지 않습니다.", 0),
        ("찰칵레시피의 답은 이렇습니다. 냉장고 사진을 찰칵 찍어 올리면:", 0, {"bold": True, "color": CORAL_DARK}),
        ("AI가 사진 속 재료를 알아봅니다 (비전 AI).", 1, {"bullet": "①"}),
        ("그 재료로 만들 수 있는 레시피 3개를 추천해줍니다 (텍스트 AI).", 1, {"bullet": "②"}),
        ("특별한 기능이 아니라, 누구나 공감하는 작은 불편에서 출발했다는 점이 중요합니다.", 0, {"bold": True}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "좋은 프로젝트는 거창한 아이디어에서 시작하지 않는다는 것을 강조. 학생들에게 "
                 "\"여러분이 오늘 아침 겪은 사소한 불편은 무엇이었나요?\"라고 질문을 던져 잠깐 대화를 "
                 "유도해도 좋다.")

    s = content_slide(prs, "Part 1 · 우리가 만든 것", "사용자 흐름 4단계", nx(), TOTAL)
    steps = [
        ("① 사진 업로드", "냉장고 안을 최대 3장까지 찍어 올립니다."),
        ("② 재료 확인", "AI가 알아낸 재료 목록을 사람이 직접 고치고 보완합니다."),
        ("③ 레시피 추천", "확정된 재료와 내 취향(맵기·난이도·시간)으로 레시피 3개를 받습니다."),
        ("④ 평가와 저장", "마음에 든 레시피는 저장하고, 좋아요·코멘트로 기록을 남깁니다."),
    ]
    box_w = Inches(2.78)
    gap = Inches(0.15)
    for i, (title, desc) in enumerate(steps):
        x = Inches(0.7) + (box_w + gap) * i
        add_rect(s, x, Inches(2.15), box_w, Inches(3.7), fill=WHITE, line_color=LINE, line_w=1.2, radius=0.08)
        add_rect(s, x, Inches(2.15), box_w, Inches(0.85), fill=CORAL, radius=0.08)
        add_text(s, title, x, Inches(2.15), box_w, Inches(0.85), size=17, color=WHITE, bold=True,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc, x + Inches(0.18), Inches(3.15), box_w - Inches(0.36), Inches(2.5),
                  size=14.5, color=INK, line_spacing=1.3)
        if i < 3:
            add_text(s, "→", x + box_w, Inches(2.15), gap + Inches(0.3), Inches(0.85), size=20,
                      color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    set_notes(s, "화면 상단에 4단계 진행 표시(마법사 스텝)가 실제로 붙어 있다는 점을 언급하면 좋다.")

    s = content_slide(prs, "Part 1 · 우리가 만든 것", "왜 이 프로젝트가 좋은 교재인가?", nx(), TOTAL)
    add_bullets(s, [
        ("실제로 배포되어 지금 이 순간에도 접속할 수 있는 살아있는 서비스입니다.", 0),
        ("회원가입·로그인, 데이터베이스, 외부 AI API, 배포까지 실무 웹서비스의 핵심 요소를 모두 담고 "
         "있습니다.", 0),
        ("그러면서도 복잡한 결제·물류 같은 요소는 없어 신입생이 전체 그림을 한 번에 이해하기 좋습니다.", 0),
        ("무엇보다, 오늘 보여드릴 것은 \"완성된 결과\"가 아니라 \"만들어져 온 과정\" 그 자체입니다.", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "이 슬라이드는 Part1의 마무리이자 Part2로 넘어가는 다리 역할이다.")

    # ================= PART 2 =================
    section_divider(prs, 2, "바이브 코딩이란 무엇인가", "AI와 대화하며 만드는 새로운 개발 방식")
    nx()

    s = content_slide(prs, "Part 2 · 바이브 코딩", "전통적 개발 vs 바이브 코딩", nx(), TOTAL)
    col_w = Inches(5.7)
    add_rect(s, Inches(0.7), Inches(2.0), col_w, Inches(4.2), fill=WHITE, line_color=LINE, line_w=1.2, radius=0.06)
    add_text(s, "전통적인 개발", Inches(0.7), Inches(2.15), col_w, Inches(0.5), size=19, bold=True,
              color=INK, align=PP_ALIGN.CENTER)
    add_bullets(s, [
        ("문법을 하나하나 배우고 외웁니다.", 0),
        ("에러 메시지를 검색하며 원인을 추적합니다.", 0),
        ("기능 하나 추가에도 관련 문서를 오래 뒤집니다.", 0),
        ("\"어떻게 구현하는가\"에 대부분의 시간을 씁니다.", 0),
    ], Inches(0.95), Inches(2.75), col_w - Inches(0.5), Inches(3.2), size=17)
    add_rect(s, Inches(6.65), Inches(2.0), col_w, Inches(4.2), fill=CORAL, radius=0.06)
    add_text(s, "바이브 코딩", Inches(6.65), Inches(2.15), col_w, Inches(0.5), size=19, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
    add_bullets(s, [
        ("무엇을 원하는지 사람의 말로 설명합니다.", 0, {"color": CREAM2}),
        ("AI가 코드를 직접 작성·실행하고 결과를 보여줍니다.", 0, {"color": CREAM2}),
        ("사람은 결과를 보고 맞다·아니다를 판단합니다.", 0, {"color": CREAM2}),
        ("\"무엇을, 왜 만드는가\"에 대부분의 시간을 씁니다.", 0, {"color": CREAM2}),
    ], Inches(6.9), Inches(2.75), col_w - Inches(0.5), Inches(3.2), size=17, color=CREAM2)
    set_notes(s, "두 방식이 대립한다기보다, 힘을 쏟는 지점이 옮겨갔다는 것을 강조한다.")

    s = content_slide(prs, "Part 2 · 바이브 코딩", "클로드 코드(Claude Code)란?", nx(), TOTAL)
    add_bullets(s, [
        ("터미널(명령창)에서 실행하는 AI 개발 동료입니다.", 0, {"bold": True}),
        ("코드를 \"읽고 → 쓰고 → 실행하고 → 결과를 확인\"하는 과정을 스스로 반복할 수 있습니다.", 0),
        ("파일을 만들고 고치고, 터미널 명령을 실행하고, 브라우저까지 직접 열어 화면을 확인합니다.", 0),
        ("사람은 방향을 정하고 결과를 검수하는 역할에 집중할 수 있습니다.", 0),
        ("오늘 이 강의 자료도, 찰칵레시피 앱도 전부 클로드 코드와의 대화로 만들어졌습니다.", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "클로드 코드를 단순 챗봇이 아니라 직접 일하는 동료로 소개하는 것이 핵심이다.")

    s = content_slide(prs, "Part 2 · 바이브 코딩", "바이브 코딩의 4단계 루프", nx(), TOTAL)
    loop = [
        ("① 설명한다", "이런 화면·기능이 필요해, 라고 사람의 말로 요청합니다."),
        ("② 구현한다", "AI가 코드를 작성하고 필요하면 직접 실행합니다."),
        ("③ 검증한다", "화면을 함께 보며 이게 맞는지 확인합니다."),
        ("④ 반복한다", "안 맞으면 다시 설명합니다. 이 순환을 계속 돕니다."),
    ]
    box_w = Inches(2.78)
    for i, (title, desc) in enumerate(loop):
        x = Inches(0.7) + (box_w + Inches(0.15)) * i
        add_rect(s, x, Inches(2.3), box_w, Inches(3.3), fill=CREAM2, radius=0.1)
        add_text(s, title, x, Inches(2.5), box_w, Inches(0.6), size=19, bold=True, color=CORAL_DARK,
                  align=PP_ALIGN.CENTER)
        add_text(s, desc, x + Inches(0.2), Inches(3.15), box_w - Inches(0.4), Inches(2.2), size=14.5,
                  color=INK, line_spacing=1.3, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(s, "↻" if i == 2 else "→", x + box_w, Inches(2.3), Inches(0.3), Inches(3.3),
                      size=20, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "④ 다음에는 다시 ①로. 이 루프를 수백 번 돌린 결과가 바로 오늘의 찰칵레시피입니다.",
              Inches(0.7), Inches(5.9), Inches(11.8), Inches(0.6), size=15, color=MUTED, italic=True,
              align=PP_ALIGN.CENTER)
    set_notes(s, "이 루프 다이어그램이 오늘 강의 전체를 관통하는 핵심 그림이다.")

    s = content_slide(prs, "Part 2 · 바이브 코딩", "오늘 실제로 나눴던 대화, 미리 살짝 보기", nx(), TOTAL)
    add_bullets(s, [
        ("관리자만 보던 레시피 평가를, 회원 누구나 볼 수 있게 바꿔줘.", 0),
        ("사진과 재료 인식 결과를 사용자와 AI가 함께 채점하게 만들고 싶어.", 0),
        ("모델 드롭다운 UI가 카드 밖으로 삐져나와.", 0),
        ("세 문장 모두 코드 용어가 하나도 없습니다. 그런데 전부 실제로 구현됐습니다.", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "실제로 이 강의를 준비하며 나눈 대화에서 그대로 가져온 문장들이라는 걸 밝힌다. 이전 판(v1)에서 "
                 "이 문장들에 따옴표 기호를 말머리처럼 붙였더니 어색하다는 피드백이 있어, 이번엔 평범한 "
                 "글머리표로 정리했다.")

    s = worksheet_slide(prs, nx(), TOTAL, "실습 — 여러분의 첫 프롬프트를 써봅시다", [
        ("좋은 프롬프트의 3요소: 무엇을, 왜, 어떤 조건으로.", "예: \"사진 업로드 버튼을 눌렀을 때 (무엇을) "
         "업로드 중이라는 걸 알 수 있게 (왜) 스피너를 보여줘 (어떤 조건으로).\""),
        ("나쁜 프롬프트 예시: \"업로드 좀 예쁘게 해줘.\"", "무엇이 예쁜 것인지 기준이 없어, AI가 임의로 "
         "판단해야 합니다."),
        ("좋은 프롬프트로 고쳐보면?", "\"업로드 버튼을 카드 안에서 가로로 꽉 차게, 우리 테마의 주황색 "
         "포인트 색으로 바꿔줘.\" 처럼 구체적인 기준을 넣어봅니다."),
        ("이제 여러분 차례입니다. 평소 쓰던 앱에서 아쉬웠던 화면 하나를 골라 프롬프트를 써보세요.", "짝과 "
         "바꿔 읽고, \"이 문장만 보고 무엇을 만들어야 할지 알겠는지\" 서로 확인해봅니다."),
    ], "프롬프트 작성 실습이다. \"무엇을, 왜, 어떤 조건으로\"라는 3요소 틀을 반복해서 강조한다. 나쁜 "
       "예시와 좋은 예시를 나란히 보여주는 것이 효과적이다. 시간이 있다면 실제로 학생들이 쓴 프롬프트 "
       "몇 개를 발표시켜도 좋다.")

    s = content_slide(prs, "Part 2 · 바이브 코딩", "AI에게 안전하게 맡기는 법", nx(), TOTAL)
    add_bullets(s, [
        ("비밀번호, API 키, 주민등록번호 같은 값은 절대 프롬프트에 그대로 붙여넣지 않습니다. 예시나 "
         "가짜 값으로 대체해서 설명합니다.", 0, {"bold": True, "color": CORAL_DARK}),
        ("결제, 회원 탈퇴, 데이터 삭제처럼 되돌리기 어려운 작업은 AI가 실행하기 전에 반드시 사람이 한 "
         "번 더 확인합니다.", 0),
        ("AI가 만든 코드도 결국 내 이름으로 배포되는 것입니다. \"AI가 짰으니까 몰라도 된다\"가 아니라, "
         "무엇이 왜 그렇게 동작하는지 마지막엔 내가 이해하고 있어야 합니다.", 0),
        ("오늘 배운 4단계 루프의 \"③ 검증한다\" 단계가 바로 이 책임을 지키는 과정입니다.", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=18)
    set_notes(s, "바이브 코딩의 편리함만 강조하면 안전 불감증으로 이어질 수 있다. 이 슬라이드에서 최소한의 "
                 "안전 수칙을 짚어주고, Part 8의 \"상식적 판단력\", \"클로드 코드 협업에 대한 신뢰\" "
                 "역량과 연결지어 설명한다.")

    # ================= 용어집 (미니 파트) =================
    s = content_slide(prs, "잠깐, 용어 정리", "오늘 나온 용어, 한 번에 정리 (1/2)", nx(), TOTAL)
    terms1 = [
        ("API", "서로 다른 프로그램이 대화하는 규칙과 창구. \"이 주소로 이렇게 물어보면 이런 답을 준다\"는 "
         "약속입니다."),
        ("서버 / 클라이언트", "서버는 데이터를 관리하고 응답하는 쪽, 클라이언트는 사용자가 보는 화면(브라우저) "
         "쪽입니다."),
        ("배포(deploy)", "내 컴퓨터에만 있던 코드를, 누구나 접속할 수 있는 인터넷 주소로 실제 서비스하는 "
         "것입니다."),
        ("환경변수", "API 키처럼 코드에 직접 적으면 안 되는 비밀 값을, 코드 밖에 따로 보관해두는 방식입니다."),
        ("데이터베이스", "회원 정보나 게시글처럼, 계속 저장해두고 나중에 다시 꺼내 써야 하는 정보를 담는 "
         "곳입니다."),
        ("컴포넌트", "화면의 한 부분(버튼, 카드, 목록 등)을 재사용 가능한 조각으로 만들어둔 것입니다."),
        ("커밋 / 푸시", "커밋은 코드 변경을 사진처럼 기록하는 것, 푸시는 그 기록을 GitHub 같은 원격 저장소에 "
         "올리는 것입니다."),
        ("RLS(행 단위 보안)", "데이터베이스의 각 표(행)마다 \"누가 볼 수 있는가\"를 정하는 보안 규칙입니다."),
    ]
    top = Inches(1.95)
    row_h = Inches(0.6)
    for i, (term, desc) in enumerate(terms1):
        y = top + row_h * i
        add_text(s, term, Inches(0.7), y, Inches(2.5), row_h, size=15, color=CORAL_DARK, bold=True,
                  anchor=MSO_ANCHOR.TOP)
        add_text(s, desc, Inches(3.3), y, Inches(9.2), row_h, size=13.5, color=INK, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.15)
    set_notes(s, "용어집은 지금까지 나온 낯선 단어들을 한 번 더 정리해주는 역할이다. 이 두 장은 강의 중 "
                 "빠르게 훑고 지나가되, 강의 자료로 남겨 나중에 학생들이 사전처럼 찾아볼 수 있게 한다.")

    s = content_slide(prs, "잠깐, 용어 정리", "오늘 나온 용어, 한 번에 정리 (2/2)", nx(), TOTAL)
    terms2 = [
        ("프롬프트", "AI에게 무엇을 해달라고 요청하는 문장이나 지시문입니다."),
        ("토큰", "AI가 글을 처리하는 최소 단위입니다. 글자 그대로가 아니라 조각 단위로 세며, 요금이나 "
         "속도 계산의 기준이 됩니다."),
        ("무료 티어(free tier)", "돈을 내지 않아도 일정 한도까지 쓸 수 있게 해주는 요금제입니다."),
        ("Fallback(대체 실행)", "1순위가 실패했을 때 자동으로 2순위, 3순위로 넘어가 다시 시도하는 "
         "안전장치입니다."),
        ("mock(모형)", "진짜 대신 잠깐 쓰는 가짜 구현입니다. 진짜가 준비되기 전에 전체 흐름을 먼저 "
         "검증할 때 씁니다."),
        ("OAuth", "\"비밀번호를 새로 만들지 않고, 구글·애플 같은 다른 서비스의 로그인을 빌려 쓰는\" "
         "인증 방식입니다."),
        ("ERD", "데이터베이스의 표들이 서로 어떻게 연결되는지 그린 설계도입니다."),
        ("PRD", "무엇을, 왜, 어떻게 만들지 코드 짜기 전에 정리해두는 요구사항 문서입니다."),
    ]
    for i, (term, desc) in enumerate(terms2):
        y = top + row_h * i
        add_text(s, term, Inches(0.7), y, Inches(2.5), row_h, size=15, color=CORAL_DARK, bold=True,
                  anchor=MSO_ANCHOR.TOP)
        add_text(s, desc, Inches(3.3), y, Inches(9.2), row_h, size=13.5, color=INK, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.15)
    set_notes(s, "두 번째 용어집 슬라이드다. 이후 Part3(PRD), Part4(사이트별 준비), Part5(기술 심화)에서 "
                 "이 단어들이 실제로 다시 등장하니, 여기서 미리 익숙해지도록 한다.")

    # ================= PART 3: PRD 작성 과정 상세 =================
    section_divider(prs, 3, "PRD 작성 과정 상세 안내", "코드를 한 줄도 쓰기 전에, 무엇을 어떻게 정리했는가")
    nx()

    s = content_slide(prs, "Part 3 · PRD", "PRD가 무엇이고, 왜 코드보다 먼저 쓰는가?", nx(), TOTAL)
    add_bullets(s, [
        ("PRD는 \"Product Requirements Document\"의 줄임말입니다. 한국어로는 \"제품 요구사항 문서\" "
         "정도로 부를 수 있습니다.", 0),
        ("쉽게 말하면 \"우리가 무엇을, 왜, 어떻게 만들 것인지\"를 코드를 짜기 전에 글로 정리한 메모입니다.", 0),
        ("집을 지을 때 설계도 없이 벽부터 세우지 않는 것처럼, 앱도 화면과 기능을 정하지 않고 코드부터 짜면 "
         "나중에 다시 뜯어고치는 일이 훨씬 많아집니다.", 0),
        ("바이브 코딩에서는 이 PRD 자체도 클로드 코드와 대화하며 함께 씁니다. 문서 작성도 대화의 일부입니다.", 0,
         {"bold": True, "color": CORAL_DARK}),
        ("찰칵레시피 저장소의 docs/PRD.md 파일이 바로 이 문서입니다. 지금도 그대로 남아 있습니다.", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.5), size=18)
    set_notes(s, "PRD라는 용어를 처음 듣는 학생이 많을 것이다. 전공 용어로 겁먹지 않도록 \"만들기 전에 쓰는 "
                 "메모\"라는 쉬운 말로 먼저 풀어주고, 이후 슬라이드에서 실제 파일 구조를 보여준다.")

    s = content_slide(prs, "Part 3 · PRD", "우리 프로젝트의 실제 PRD 구조 — 10개 장", nx(), TOTAL)
    toc = [
        ("1장", "개요", "한 문단으로 이 앱이 무엇인지 요약"),
        ("2장", "목표와 범위", "무엇을 할 것인가, 그리고 무엇은 하지 않을 것인가"),
        ("3장", "기술 스택", "어떤 도구·서비스로 만들 것인가"),
        ("4장", "사용자 흐름", "화면이 어떤 순서로 이어지는가"),
        ("5장", "기능 요구사항", "각 화면이 구체적으로 무엇을 해야 하는가"),
        ("6장", "데이터 모델", "무엇을, 어떤 표 구조로 저장할 것인가"),
        ("7장", "API/AI 연동 설계", "외부 서비스를 어떻게, 어떤 순서로 연결하는가"),
        ("8장", "비기능 요구사항", "보안·성능·접근성처럼 눈에 안 보이지만 중요한 것"),
        ("9장", "마일스톤", "몇 단계로 나누어 만들 것인가"),
        ("10장", "미결정 사항", "아직 정하지 못한 것을 정직하게 적어두는 곳"),
    ]
    col_items = [toc[:5], toc[5:]]
    for c, items in enumerate(col_items):
        x = Inches(0.7) + Inches(6.05) * c
        for i, (num, name, desc) in enumerate(items):
            y = Inches(2.0) + Inches(0.86) * i
            add_rect(s, x, y, Inches(5.8), Inches(0.78), fill=CREAM2 if i % 2 == 0 else
                     RGBColor(0xFA, 0xF3, 0xE8), radius=0.1)
            add_text(s, num, x + Inches(0.15), y, Inches(0.8), Inches(0.78), size=15, color=CORAL_DARK,
                      bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, name, x + Inches(0.95), y + Pt(6), Inches(2.1), Inches(0.35), size=14.5,
                      color=INK, bold=True)
            add_text(s, desc, x + Inches(0.95), y + Inches(0.38), Inches(4.7), Inches(0.35), size=11.5,
                      color=MUTED)
    set_notes(s, "이 10개 장은 어느 프로젝트에도 그대로 옮겨 쓸 수 있는 뼈대다. 학생들에게 \"여러분의 "
                 "아이디어도 이 10칸을 채우면 PRD가 완성된다\"고 알려준다. 이어지는 슬라이드에서 각 장을 "
                 "실제 예시와 함께 하나씩 풀어본다.")

    s = content_slide(prs, "Part 3 · PRD", "1~3장 — 무엇을, 왜, 무엇으로 만들 것인가", nx(), TOTAL)
    add_bullets(s, [
        ("1장 개요의 실제 문장: \"냉장고 사진을 업로드하면 AI가 식재료를 인식하고, 사용자의 선호 조건을 "
         "반영해 레시피를 추천하는 웹 앱.\"", 0),
        ("2장 목표와 범위는 표 하나로 정리했습니다. 목표는 \"오늘 뭐 해먹지에 대한 실용적 답을 빠르게 "
         "제공\", 비목표는 \"칼로리 계산, 재고 관리, 모바일 네이티브 앱\"이었습니다.", 0),
        ("비목표를 적는 것이 의외로 중요합니다. 무엇을 안 할지 정해야, 무엇을 할지도 명확해집니다.", 0,
         {"bold": True}),
        ("3장 기술 스택에서는 Next.js, Vercel, Supabase, OpenRouter처럼 사용할 도구 이름을 미리 정해 "
         "뒀습니다. (이 도구들의 역할은 Part 4에서 자세히 다룹니다.)", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.5), size=17.5)
    set_notes(s, "실제 PRD.md의 문장을 그대로 인용해 \"이 정도 수준으로만 쓰면 된다\"는 감을 잡아준다. "
                 "너무 길거나 학술적으로 쓸 필요가 없다는 걸 알려준다.")

    s = content_slide(prs, "Part 3 · PRD", "4장 — 사용자 흐름을 그림 없이 글로 적는 법", nx(), TOTAL)
    add_bullets(s, [
        ("화면 디자인을 그리기 전에, 화면이 이어지는 순서를 번호가 매겨진 문장으로 먼저 적습니다.", 0),
        ("실제 PRD 4장은 이렇게 시작합니다: \"1. 가입/로그인 → 2. 선호 조건 설정 → 3. 냉장고 사진 업로드 "
         "→ 4. 식재료 인식 결과 확인 → 5. 레시피 추천 요청 → 6. 레시피 상세 확인, 평가, 저장\"", 0),
        ("각 단계마다 \"이 화면에서 사용자가 무엇을 볼 수 있고, 무엇을 할 수 있는가\"를 한두 문장으로 "
         "덧붙입니다.", 0),
        ("이 순서가 그대로 화면 상단의 마법사 스텝 표시(1. 사진 업로드 → 2. 재료 확인 → 3. 레시피 추천)로 "
         "구현됐습니다. PRD의 문장이 실제 화면 구조가 된 것입니다.", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.5), size=17.5)
    set_notes(s, "PRD의 문장이 실제 UI로 그대로 이어졌다는 것을 보여주는 것이 핵심이다. 기획 문서가 추상적인 "
                 "낙서가 아니라 실제 구현의 설계도라는 것을 체감시킨다.")

    s = content_slide(prs, "Part 3 · PRD", "5~6장 — 해야 할 일과 저장할 정보를 목록·표로", nx(), TOTAL)
    add_bullets(s, [
        ("5장 기능 요구사항은 화면별로 나눠 적습니다. 예: \"5.3 식재료 인식\" 아래에 \"이미지 최대 3장 "
         "업로드\", \"인식 결과에 대한 수동 추가/수정/삭제 UI\" 같은 항목을 불릿으로 나열합니다.", 0),
        ("6장 데이터 모델은 아직 데이터베이스를 만들기 전에, \"어떤 표(테이블)가 필요하고 그 표에 어떤 "
         "칸(컬럼)이 있을지\"를 미리 글로 그려보는 곳입니다.", 0),
        ("실제 예: fridge_sessions 표에는 id, user_id, vision_provider, created_at 같은 칸이 필요하다고 "
         "미리 적어뒀습니다.", 0),
        ("이렇게 미리 적어두면, 나중에 Supabase에서 실제 표를 만들 때 그대로 옮기기만 하면 됩니다.", 0,
         {"bold": True}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.5), size=17.5)
    set_notes(s, "6장 데이터 모델 파트는 나중에 Part 5의 ERD 다이어그램과 정확히 연결된다는 것을 예고하면 "
                 "좋다.")

    s = content_slide(prs, "Part 3 · PRD", "7~10장 — 연결 설계, 비기능, 일정, 그리고 솔직한 미결정", nx(), TOTAL)
    add_bullets(s, [
        ("7장 API/AI 연동 설계에서는 \"비전 API가 실패하면 다른 모델로 자동 재시도한다\"처럼 예외 상황까지 "
         "미리 정해뒀습니다.", 0),
        ("8장 비기능 요구사항은 보안(로그인한 사용자만 자기 데이터를 본다), 성능(무료 API 한도 안에서 "
         "동작), 접근성(모바일 화면 대응)처럼 화면에 직접 보이지 않지만 반드시 지켜야 할 약속입니다.", 0),
        ("9장 마일스톤은 \"1. PRD 확정 → 2. 디자인 시안 → 3. MVP → … → 8. 문서화\"처럼 큰 단계를 "
         "순서대로 나눈 것입니다.", 0),
        ("10장 미결정 사항이 어쩌면 가장 중요합니다. \"아직 정하지 못했다\"고 정직하게 적어두면, 나중에 "
         "그 결정을 놓치지 않고 다시 꺼내 대화할 수 있습니다.", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.5), size=17.5)
    set_notes(s, "10장의 가치를 특히 강조한다. 실제로 이 프로젝트에서 10장에 적어둔 \"아직 남음\" 항목들이 "
                 "나중에 Part 5·6에서 실제로 하나씩 해결되는 과정을 보여줄 것이라고 예고한다.")

    s = code_slide(
        prs, nx(), TOTAL, "Part 3 · PRD", "실제로 이렇게 대화하며 PRD를 써 내려갔습니다",
        "docs/PRD.md (실제 저장소 파일)",
        [
            "나:  냉장고 사진으로 레시피 추천하는 앱을 만들고 싶어.",
            "     사진 업로드 → 재료 인식 → 레시피 추천 흐름으로.",
            "",
            "클로드 코드:  네, PRD 초안을 먼저 정리하겠습니다.",
            "  이런 질문에 답을 주시면 문서로 정리할게요:",
            "  1) 사진은 몇 장까지 받을까요?",
            "  2) 레시피는 몇 개씩 추천할까요?",
            "  3) 회원가입은 어떤 방식으로 할까요?",
            "",
            "나:  3장까지, 3개씩, 구글 로그인이랑 이메일 둘 다.",
            "",
            "클로드 코드:  docs/PRD.md 파일을 만들었습니다. 이대로",
            "  진행해도 될까요?",
        ],
        "실제 대화는 이보다 더 짧고 자연스러운 경우도 많습니다. 핵심은 \"완벽한 질문\"이 아니라 "
        "\"질문과 답이 오가며 문서가 점점 정확해지는 과정\" 그 자체입니다.",
        "이 대화 예시는 실제 첫 대화를 단순화해 재구성한 것이다. 학생들에게 \"이 정도 대화면 충분히 "
        "시작할 수 있다\"는 안도감을 준다. 완벽한 초기 요구사항이 없어도 대화로 다듬어간다는 것을 "
        "강조한다.",
    )

    s = worksheet_slide(prs, nx(), TOTAL, "여러분의 PRD 초안 워크시트 — 지금 답해봅시다", [
        ("여러분이 겪는 작은 불편은 무엇인가요?", "예: 자취방 냉장고 재고를 자꾸 깜빡한다."),
        ("그 앱은 무엇을 하고, 무엇은 하지 않나요? (목표/비목표)", "예: 재고 알림은 하되, 자동 주문은 하지 않는다."),
        ("화면은 어떤 순서로 이어지나요? (사용자 흐름)", "예: 로그인 → 재고 입력 → 알림 설정 → 알림 확인."),
        ("어떤 정보를 저장해야 하나요? (데이터 모델)", "예: 사용자, 재고 항목, 유통기한, 알림 설정."),
        ("아직 정하지 못한 것은 무엇인가요? (미결정 사항)", "예: 유통기한을 사진으로 인식할지, 직접 입력할지."),
    ], "이 워크시트를 강의 중 짝과 함께 채워보게 하거나, 다음 시간까지 과제로 내줘도 좋다. 다섯 질문에 "
       "한두 문장씩만 답해도 PRD 1장 분량의 뼈대가 완성된다는 것을 체감시킨다.")

    s = content_slide(prs, "Part 3 · PRD", "PRD는 한 번 쓰고 끝이 아닙니다", nx(), TOTAL)
    add_bullets(s, [
        ("찰칵레시피 저장소의 docs/ 폴더에는 PRD00부터 PRD06까지, 총 7개의 스냅샷이 남아 있습니다.", 0),
        ("PRD00(최초 초안) → PRD01(디자인 착수 전) → PRD02(디자인 확정 후) → PRD03(데이터 모델링 완료) "
         "→ PRD04(MVP 완성) → PRD05(실제 API 연동·배포) → PRD06(현재 상태).", 0),
        ("프로젝트가 진행되며 계획이 바뀌거나 새로운 기능이 생길 때마다, PRD 문서도 함께 갱신했습니다.", 0),
        ("문서를 살아있는 기록으로 유지하는 것도 개발의 일부입니다. 이 습관이 나중에 \"우리가 왜 이렇게 "
         "만들었더라?\"를 되짚어볼 때 큰 힘이 됩니다.", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.4), size=18)
    set_notes(s, "PRD가 정적인 문서가 아니라 프로젝트와 함께 자라는 살아있는 기록이라는 것을 강조한다. "
                 "Part 6의 문서화 사례와 자연스럽게 연결된다.")

    # ================= PART 4: 5개 사이트 실습 가이드 =================
    section_divider(prs, 4, "개발 착수 전 준비", "5개 사이트에 계정을 만들고 연결하기 — GitHub 포함")
    nx()

    s = content_slide(prs, "Part 4 · 계정 준비", "왜 코드를 짜기 전에 계정부터 만드는가?", nx(), TOTAL)
    add_bullets(s, [
        ("찰칵레시피는 혼자 다 만든 것이 아닙니다. 5개의 외부 서비스와 함께 만들었습니다.", 0),
        ("GitHub(코드 보관·이력 관리), Vercel(배포), Supabase(로그인·데이터베이스·사진 저장), "
         "OpenRouter(AI 모델 호출), Google Cloud Console(구글 로그인 열쇠 발급)입니다.", 0),
        ("다섯 곳 모두 무료로 시작할 수 있습니다. 오늘 이 순서대로 하나씩 가입하고 연결해보겠습니다.", 0,
         {"bold": True}),
        ("각 사이트는 \"가입 → 새 프로젝트 만들기 → 열쇠(키) 하나 발급받기\"라는 같은 패턴을 반복한다는 "
         "것을 기억하면 낯설지 않습니다.", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=18)
    set_notes(s, "Part4 전체의 문을 여는 슬라이드. 다섯 사이트를 한 번에 나열해 부담스러워 보일 수 있지만, "
                 "공통 패턴이 있다는 것을 미리 알려줘 안심시킨다.")

    s = content_slide(prs, "Part 4 · 계정 준비", "왜 하필 이 순서인가?", nx(), TOTAL)
    add_bullets(s, [
        ("GitHub을 가장 먼저 준비합니다. 코드를 담을 그릇이 있어야, 그 다음 단계인 배포도 의미가 "
         "있습니다.", 0),
        ("Vercel은 GitHub 저장소가 있어야 연결할 수 있으므로 두 번째입니다.", 0),
        ("Supabase와 OpenRouter, Google Cloud Console은 서로 순서를 바꿔도 무방합니다. 다만 로그인 "
         "기능을 쓰려면 Google Cloud Console에서 발급받은 키를 Supabase에 등록해야 하므로, 두 "
         "사이트는 서로 연결되어 있다는 것을 기억해야 합니다.", 0),
        ("즉 \"코드를 담을 그릇 → 그 그릇을 세상에 공개하는 곳 → 그 안에서 쓸 부품들\" 순서로 이해하면 "
         "됩니다.", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=18)
    set_notes(s, "다섯 사이트를 그냥 순서대로 나열한 것이 아니라, 서로 의존하는 관계가 있다는 것을 "
                 "이해시켜준다. 이 감각이 있어야 나중에 다른 프로젝트를 준비할 때도 스스로 순서를 판단할 "
                 "수 있다.")

    # ---- GitHub ----
    s = site_slide(
        prs, nx(), TOTAL, "Part 4 · 계정 준비 ①", "github.com",
        "코드의 사진첩이자, 다른 모든 서비스를 잇는 다리",
        "내가 만든 코드의 모든 순간을 사진처럼 저장해두고, 필요하면 언제든 그 순간으로 되돌아갈 수 있는 곳입니다.",
        [
            ("코드를 저장하고, 바뀐 이력을 전부 기록합니다. 이걸 \"버전 관리\"라고 부릅니다.", 0),
            ("Vercel이 배포할 코드를 가져오는 곳도 바로 여기입니다. 즉 GitHub은 다른 서비스들을 잇는 "
             "다리 역할을 합니다.", 0),
            ("클로드 코드는 이 저장소에 직접 커밋(변경 기록)을 남기고 푸시(업로드)할 수 있습니다.", 0),
            ("오늘 우리가 볼 코드 스니펫, 커밋 메시지, 배포 이력이 전부 이 GitHub 저장소 안에 남아 있습니다.", 0,
             {"bold": True}),
        ],
        "GitHub을 5개 사이트 중 가장 먼저 소개하는 이유는, 나머지 대부분(특히 Vercel)이 GitHub 저장소를 "
        "전제로 동작하기 때문이다. \"모든 것의 출발점\"이라는 위치를 분명히 해준다.",
        NAVY,
    )

    s = content_slide(prs, "Part 4 · 계정 준비 ①", "버전 관리 — 왜 \"사진첩\"에 비유하는가?", nx(), TOTAL)
    add_bullets(s, [
        ("코드를 고칠 때마다 파일 이름을 \"app_final.js\", \"app_final2.js\"처럼 바꿔가며 백업해본 적이 "
         "있나요? 버전 관리는 이걸 자동으로, 훨씬 안전하게 해줍니다.", 0),
        ("\"커밋(commit)\"은 지금 이 순간의 코드 전체를 사진처럼 찍어 남기는 것입니다.", 0),
        ("사진마다 \"드롭다운 오버플로 버그 수정\"처럼 설명을 붙여두면, 나중에 무엇을 언제 왜 고쳤는지 "
         "한눈에 볼 수 있습니다.", 0),
        ("문제가 생기면 이전 \"사진\"으로 언제든 되돌아갈 수 있습니다. 이것이 Part 7에서 다룰 \"실패해도 "
         "괜찮은 이유\"의 진짜 근거입니다.", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=18)
    set_notes(s, "버전 관리라는 개념 자체를 처음 듣는 학생들을 위한 슬라이드다. 게임의 \"세이브 포인트\"에 "
                 "비유해도 좋다.")

    s = site_step_slide(
        prs, nx(), TOTAL, "Part 4 · 계정 준비 ①", "github.com — 따라 하기",
        "저장소 하나를 처음 만드는 데 5분이면 충분합니다.",
        [
            ("계정 만들기", "github.com에 접속해 Sign up을 누르고, 이메일과 비밀번호로 가입합니다. "
             "구글·애플 계정으로도 가입할 수 있습니다."),
            ("새 저장소(repository) 만들기", "우측 상단 + 버튼 → New repository. 이름을 정하고(예: "
             "my-first-app), Public(공개) 또는 Private(비공개)을 선택한 뒤 Create repository를 누릅니다."),
            ("클로드 코드와 연결하기", "내 컴퓨터에서 프로젝트 폴더를 열고, 클로드 코드에게 \"이 프로젝트를 "
             "방금 만든 GitHub 저장소와 연결해줘\"라고 요청하면 필요한 명령을 대신 실행해줍니다."),
            ("변경사항을 저장하는 습관 들이기", "기능을 하나 완성할 때마다 \"커밋하고 푸시해줘\"라고 "
             "요청하세요. 자동으로 오늘 배운 것과 같은 기록이 차곡차곡 쌓입니다."),
        ],
        ("이제 여러분에게 이런 것이 생겼습니다", "인터넷 어디서든 열어볼 수 있는 나만의 코드 저장소 하나와, "
         "그 코드의 모든 변경 이력."),
        "실습 시간이 있다면 이 슬라이드에서 실제로 학생들이 github.com에 접속해 계정을 만들어보게 한다. "
        "GitHub 계정은 이후 Vercel 로그인에도 그대로 재사용된다는 것을 미리 알려준다.",
    )

    s = code_slide(
        prs, nx(), TOTAL, "Part 4 · 계정 준비 ①", "실제로 이렇게 커밋되었습니다",
        "git log --oneline (찰칵레시피 저장소 실제 이력 일부)",
        [
            "18fe630  Fix dropdown overflow, unify model names, add photo zoom",
            "a1afe17  Split ratings into recipe popularity vs AI model quality pages",
            "a53eee9  Rebrand, public recipe ratings, and AI-driven model quality scoring",
            "c2b9bae  Add session history for re-recommending recipes, overhaul design system",
            "ccc0bb8  Add OpenRouter vision/text providers as selectable real API option",
            "c59fb4e  Initial commit: recipe4fridge_pic MVP",
        ],
        "커밋 메시지 하나하나가 \"이 순간에 무엇이, 왜 바뀌었는지\"를 설명합니다. 이 목록만 훑어봐도 "
        "프로젝트가 어떤 순서로 자랐는지 알 수 있습니다. 클로드 코드가 매번 이런 메시지를 직접 작성해 "
        "커밋합니다.",
        "실제 저장소의 진짜 커밋 로그를 그대로 가져온 것이다. 학생들에게 \"이 목록 자체가 이 프로젝트의 "
        "일기장\"이라고 설명하면 좋다.",
    )

    s = content_slide(prs, "Part 4 · 계정 준비 ①", "GitHub — 자주 하는 실수", nx(), TOTAL)
    add_bullets(s, [
        ("Private(비공개)로 만들어야 할 저장소를 실수로 Public(공개)으로 만드는 경우가 있습니다. API 키 "
         "같은 비밀 값은 절대 코드에 직접 적지 않아야 하는 이유이기도 합니다.", 0),
        ("커밋 메시지를 비워두거나 \"수정함\"처럼 의미 없이 적으면, 나중에 이력을 봐도 무엇이 바뀌었는지 "
         "알 수 없습니다.", 0),
        ("변경 사항을 며칠씩 커밋하지 않고 쌓아두면, 나중에 무엇이 왜 바뀌었는지 되짚기 어려워집니다. "
         "작은 단위로 자주 커밋하는 것이 좋습니다.", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.0), size=17.5)
    set_notes(s, "\"자주 하는 실수\" 슬라이드는 실습 중 학생들이 실제로 자주 겪는 문제를 미리 예방하는 "
                 "역할을 한다. 시간이 부족하면 빠르게 훑고 지나가도 좋다.")

    # ---- Vercel ----
    s = site_slide(
        prs, nx(), TOTAL, "Part 4 · 계정 준비 ②", "vercel.com",
        "이 앱이 사는 집 — 배포(Deployment)",
        "내가 만든 코드를, 전 세계 누구나 접속할 수 있는 인터넷 주소로 바꿔주는 곳입니다.",
        [
            ("GitHub에 코드를 올리면, Vercel이 자동으로 감지해서 새로 빌드하고 배포합니다.", 0),
            ("그 결과가 recipe4fridge-pic.vercel.app 같은 실제 주소로 나옵니다.", 0),
            ("무료 요금제로 개인 프로젝트를 시작하기에 충분합니다.", 0),
            ("우리는 \"배포\"라는 어려운 단어 대신, \"저장 버튼을 누르면 자동으로 세상에 공개된다\"고 "
             "기억하면 됩니다.", 0, {"bold": True}),
        ],
        "Vercel을 \"내 코드가 실제로 사는 집\"이라고 소개한다. GitHub push 한 번에 자동 배포되는 흐름을 "
        "\"편지를 우체통에 넣으면 자동으로 상대방 집까지 배달되는 것\"처럼 설명해도 좋다.",
        CORAL,
    )

    s = site_step_slide(
        prs, nx(), TOTAL, "Part 4 · 계정 준비 ②", "vercel.com — 따라 하기",
        "GitHub 계정이 있다면 배포까지 채 5분이 걸리지 않습니다.",
        [
            ("가입하기", "vercel.com에서 Sign Up을 누르고 \"Continue with GitHub\"를 선택합니다. GitHub "
             "계정으로 바로 연동 가입됩니다."),
            ("새 프로젝트 만들기", "Add New → Project → 방금 만든 GitHub 저장소를 목록에서 찾아 Import를 "
             "누릅니다."),
            ("환경변수 설정하기", "Settings → Environment Variables 메뉴에서 OPENROUTER_API_KEY, "
             "SUPABASE_URL 같은 비밀 값을 등록합니다. 이 값들은 코드에 직접 적지 않습니다."),
            ("배포 버튼 누르기", "Deploy를 누르면 1~2분 안에 xxx.vercel.app 주소가 발급됩니다. 그 주소로 "
             "직접 접속해 확인합니다."),
            ("그 다음부터는 자동", "이후로는 GitHub에 코드를 push할 때마다 Vercel이 알아서 새로 배포합니다. "
             "매번 이 과정을 반복할 필요가 없습니다."),
        ],
        ("이제 여러분에게 이런 것이 생겼습니다", "전 세계 어디서든 접속 가능한 실제 웹 주소 하나와, "
         "코드를 올리기만 하면 자동으로 갱신되는 배포 파이프라인."),
        "환경변수(Environment Variables) 개념을 이 시점에서 짚어준다. \"비밀번호를 코드에 직접 적지 않고 "
        "금고에 따로 보관하는 것\"에 비유하면 이해가 쉽다.",
    )

    s = content_slide(prs, "Part 4 · 계정 준비 ②", "Vercel — 자주 하는 실수", nx(), TOTAL)
    add_bullets(s, [
        ("환경변수를 등록하고 나서 재배포를 하지 않으면, 예전 설정 그대로 동작합니다. 환경변수를 바꿨다면 "
         "다시 배포해야 반영됩니다.", 0),
        ("배포가 실패했을 때 화면만 보고 포기하지 말고, Vercel의 Deployment 화면에서 로그(Logs)를 열어 "
         "정확히 어느 줄에서 실패했는지 확인해야 합니다.", 0),
        ("무료 요금제의 서버리스 함수 실행 시간 제한을 넘기면 오류가 납니다. AI 호출처럼 오래 걸리는 "
         "작업은 maxDuration 설정을 늘려야 할 수 있습니다.", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.0), size=17.5)
    set_notes(s, "\"환경변수 변경 후 재배포\"는 실습 중 가장 흔한 막힘 포인트 중 하나다. 반드시 강조한다.")

    # ---- Supabase ----
    s = site_slide(
        prs, nx(), TOTAL, "Part 4 · 계정 준비 ③", "supabase.com",
        "이 앱의 기억과 창고 — 회원·데이터베이스·사진 저장",
        "나만의 작은 서버실을, 코드 한 줄 없이 클릭 몇 번으로 공짜로 빌리는 곳입니다.",
        [
            ("회원가입/로그인(Auth) — 이메일·비밀번호, 구글 로그인을 대신 관리해줍니다.", 0),
            ("데이터베이스(Postgres) — 사용자, 레시피, 평가 같은 표(table)들을 저장합니다.", 0),
            ("파일 저장소(Storage) — 업로드한 냉장고 사진을 안전하게 보관합니다.", 0),
            ("\"누가 어떤 데이터를 볼 수 있는가\"를 정하는 보안 규칙(RLS)도 여기서 설정합니다.", 0),
        ],
        "Supabase는 이 앱에서 가장 자주 등장하는 서비스다. 로그인, 데이터, 사진까지 세 가지 역할을 한 "
        "곳에서 다 한다는 것을 강조한다.",
        NAVY,
    )

    s = site_step_slide(
        prs, nx(), TOTAL, "Part 4 · 계정 준비 ③", "supabase.com — 따라 하기",
        "데이터베이스 서버실 하나를 통째로 빌리는 과정입니다.",
        [
            ("가입 및 새 프로젝트 만들기", "supabase.com에서 GitHub 계정으로 가입 후 New Project. 프로젝트 "
             "이름, 데이터베이스 비밀번호, 서버 지역(리전)을 설정합니다."),
            ("데이터베이스 표 만들기", "왼쪽 메뉴의 SQL Editor에 미리 준비한 스키마(표 생성 명령)를 붙여넣고 "
             "Run을 누릅니다. 찰칵레시피는 이 스키마를 supabase/ 폴더에 파일로 보관해뒀습니다."),
            ("API 키 확인하기", "Settings → API 메뉴에서 Project URL과 anon public key를 복사해, Vercel의 "
             "환경변수에 붙여넣습니다."),
            ("로그인 방식 켜기", "Authentication → Providers에서 Email, Google 같은 로그인 방식을 "
             "활성화합니다."),
            ("보안 규칙(RLS) 켜기", "각 표마다 Row Level Security를 켜고, \"본인 데이터만 조회 가능\" 같은 "
             "규칙(정책)을 작성합니다. 이 부분을 빼먹으면 누구나 남의 데이터를 볼 수 있게 됩니다."),
        ],
        ("이제 여러분에게 이런 것이 생겼습니다", "로그인 기능, 데이터베이스, 사진 저장소를 갖춘 나만의 "
         "백엔드 서버실 하나. 서버를 직접 설치하지 않고도요."),
        "RLS(Row Level Security)는 신입생에게 다소 어려운 개념일 수 있다. \"각 서랍마다 자물쇠를 걸고, "
        "누구 열쇠로 어떤 서랍을 열 수 있는지 정하는 것\"에 비유하면 좋다.",
    )

    s = content_slide(prs, "Part 4 · 계정 준비 ③", "Supabase — 자주 하는 실수", nx(), TOTAL)
    add_bullets(s, [
        ("표를 만들고 RLS를 켜는 것을 잊으면, 기본적으로 아무도 데이터를 읽지 못하거나(막혀서 오류가 "
         "나거나), 반대로 RLS 자체를 켜지 않아 누구나 다 볼 수 있는 상태가 됩니다. 둘 다 위험합니다.", 0),
        ("anon key(공개용 키)와 service role key(관리자용 키)를 혼동해 잘못된 곳에 쓰면 보안 사고로 "
         "이어질 수 있습니다. 클라이언트(브라우저)에는 반드시 anon key만 노출해야 합니다.", 0),
        ("데이터베이스 비밀번호를 잊어버리면 처음부터 다시 설정해야 하니, 발급 즉시 안전한 곳에 "
         "기록해둡니다.", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.0), size=17.5)
    set_notes(s, "anon key와 service role key의 차이는 보안상 매우 중요하다. 짧게라도 반드시 짚어줘야 "
                 "하는 대목이다.")

    # ---- OpenRouter ----
    s = site_slide(
        prs, nx(), TOTAL, "Part 4 · 계정 준비 ④", "openrouter.ai",
        "이 앱의 AI 두뇌 렌탈 서비스",
        "구글, 메타, 알리바바 같은 여러 회사의 AI 모델을 한 창구에서 골라 빌려 쓰는 곳입니다.",
        [
            ("사진을 보고 재료를 알아내는 \"눈\"(비전 모델)을 빌려옵니다.", 0),
            ("재료로 레시피를 써내는 \"글솜씨\"(텍스트 모델)를 빌려옵니다.", 0),
            ("무료(:free) 모델도 제공해서, 비용 걱정 없이 배우고 실험할 수 있습니다.", 0),
            ("한 모델이 응답을 못 하면 다른 모델로 자동으로 갈아타는 안전장치도 여기서 만듭니다.", 0),
        ],
        "OpenRouter를 \"AI 인력 파견 업체\"에 비유해도 좋다. 특히 무료 모델이 존재한다는 사실이 신입생들에게 "
        "\"나도 지금 당장 시작할 수 있다\"는 실질적 자신감을 준다.",
        CORAL_DARK,
    )

    s = site_step_slide(
        prs, nx(), TOTAL, "Part 4 · 계정 준비 ④", "openrouter.ai — 따라 하기",
        "무료 모델만 쓴다면 비용 등록 없이 바로 시작할 수 있습니다.",
        [
            ("가입하기", "openrouter.ai에서 Sign In을 누르고 구글 계정 등으로 간편하게 가입합니다."),
            ("API 키 발급하기", "우측 상단 프로필 → Keys → Create Key. 키에 이름을 붙이고 생성한 뒤, 그 "
             "키 값을 안전한 곳에 복사해둡니다. 이 값은 다시 볼 수 없는 경우가 많으니 바로 저장합니다."),
            ("무료 모델 찾아보기", "Models 메뉴에서 이름 끝에 :free가 붙은 모델들을 확인합니다. 찰칵레시피는 "
             "google/gemma, nvidia/nemotron 같은 무료 모델들을 사용합니다."),
            ("환경변수에 등록하기", "발급받은 키를 OPENROUTER_API_KEY라는 이름으로 프로젝트의 .env 파일과 "
             "Vercel 환경변수 양쪽에 등록합니다."),
            ("사용량 확인 습관", "Activity 메뉴에서 호출 횟수와 한도를 가끔 확인합니다. 무료 모델도 시간당 "
             "호출 한도가 있을 수 있습니다."),
        ],
        ("이제 여러분에게 이런 것이 생겼습니다", "여러 회사의 AI 모델을 코드 몇 줄로 호출할 수 있는 열쇠 "
         "하나. 비용 걱정 없이 실험할 수 있는 무료 모델 목록도 함께요."),
        "API 키를 다시 볼 수 없다는 점(발급 즉시 복사해야 한다는 점)은 실습 중 흔한 실수 포인트이니 "
        "명확히 짚어준다.",
    )

    s = content_slide(prs, "Part 4 · 계정 준비 ④", "OpenRouter — 자주 하는 실수", nx(), TOTAL)
    add_bullets(s, [
        (":free가 붙지 않은 유료 모델을 실수로 선택하면 비용이 청구될 수 있습니다. 학습 단계에서는 "
         "반드시 :free 모델만 사용합니다.", 0),
        ("인기 있는 무료 모델일수록 사용자가 몰려 응답이 불안정해질 수 있습니다. 이 프로젝트가 여러 "
         "모델로 자동 전환(fallback)하는 안전장치를 만든 이유이기도 합니다.", 0),
        ("API 키를 GitHub 코드에 실수로 커밋해버리면 노출됩니다. 반드시 .env 파일에 넣고, .gitignore에 "
         ".env가 포함되어 있는지 확인해야 합니다.", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.0), size=17.5)
    set_notes(s, "API 키가 GitHub에 실수로 올라가는 사고는 실무에서도 매우 흔하다. .gitignore 개념을 "
                 "이 시점에 짧게 짚어주면 좋다.")

    # ---- Google Cloud Console ----
    s = site_slide(
        prs, nx(), TOTAL, "Part 4 · 계정 준비 ⑤", "console.cloud.google.com",
        "구글 로그인을 위한 신분증 발급소",
        "우리 앱이 \"구글 로그인\" 버튼을 달 수 있게, 구글이 발급해주는 열쇠(OAuth) 창구입니다.",
        [
            ("여기서 \"이 앱은 이런 이름과 주소를 가진 진짜 프로젝트\"라고 등록합니다.", 0),
            ("등록하면 클라이언트 ID·비밀 키를 발급받아, 우리 앱에 \"구글로 로그인\" 버튼을 연결할 수 "
             "있습니다.", 0),
            ("사용자 입장에서는 비밀번호를 새로 만들지 않아도 되니 훨씬 편해집니다.", 0),
            ("다섯 사이트 중 가장 낯선 이름이지만, 역할은 \"신분증 발급\" 하나로 단순합니다.", 0, {"bold": True}),
        ],
        "Google Cloud Console은 이름이 가장 위압적이지만 우리가 쓰는 역할은 딱 하나, OAuth(구글 로그인) "
        "열쇠 발급이라는 것을 분명히 한다.",
        GREEN,
    )

    s = site_step_slide(
        prs, nx(), TOTAL, "Part 4 · 계정 준비 ⑤", "console.cloud.google.com — 따라 하기",
        "구글 로그인 열쇠 하나를 발급받아 Supabase에 꽂아주는 과정입니다.",
        [
            ("프로젝트 만들기", "console.cloud.google.com 접속 후 새 프로젝트를 만들고 이름을 지정합니다."),
            ("OAuth 동의 화면 구성하기", "API 및 서비스 → OAuth 동의 화면. User Type을 \"외부\"로 선택하고, "
             "앱 이름과 연락 이메일 같은 기본 정보를 입력합니다."),
            ("사용자 인증 정보 만들기", "사용자 인증 정보 → 만들기 → OAuth 클라이언트 ID. 애플리케이션 "
             "유형은 \"웹 애플리케이션\"을 선택합니다."),
            ("승인된 리디렉션 URI 등록하기", "Supabase Authentication 설정 화면이 알려주는 콜백 주소(예: "
             "https://xxxx.supabase.co/auth/v1/callback)를 여기에 그대로 등록합니다."),
            ("발급된 키를 Supabase에 연결하기", "생성된 클라이언트 ID와 보안 비밀을 Supabase의 "
             "Authentication → Providers → Google 설정 화면에 붙여넣고 저장합니다."),
        ],
        ("이제 여러분에게 이런 것이 생겼습니다", "\"구글로 로그인\" 버튼이 실제로 동작하는, 진짜 인증된 "
         "앱 하나."),
        "이 부분은 다섯 사이트 중 가장 절차가 복잡하다. 리디렉션 URI를 정확히 복사-붙여넣기 하지 않으면 "
        "흔히 오류가 난다는 것을 미리 알려주면 실습 중 혼란을 줄일 수 있다.")

    s = content_slide(prs, "Part 4 · 계정 준비 ⑤", "Google Cloud Console — 자주 하는 실수", nx(), TOTAL)
    add_bullets(s, [
        ("리디렉션 URI를 한 글자라도 다르게 입력하면(예: 끝에 슬래시 유무) \"redirect_uri_mismatch\" "
         "오류가 납니다. Supabase가 알려주는 주소를 그대로 복사해 붙여넣어야 합니다.", 0),
        ("OAuth 동의 화면을 \"테스트\" 상태로 두면, 미리 등록한 테스트 사용자 외에는 로그인이 막힐 수 "
         "있습니다. 실제 서비스로 공개하려면 \"게시\" 절차가 필요합니다.", 0),
        ("클라이언트 ID와 보안 비밀(secret)을 뒤바꿔 붙여넣는 실수도 자주 일어납니다. 두 값의 이름을 "
         "다시 한번 확인합니다.", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.0), size=17.5)
    set_notes(s, "구글 OAuth 설정은 다섯 사이트 중 가장 오류가 잦은 구간이다. 이 슬라이드에서 시간을 "
                 "조금 더 써도 좋다.")

    s = content_slide(prs, "Part 4 · 계정 준비", "환경변수 총정리 — 어디에 무엇을 넣었나", nx(), TOTAL)
    env_rows = [
        ("이름", "발급받은 곳", "어디에 등록하나"),
        ("NEXT_PUBLIC_SUPABASE_URL", "supabase.com → Settings → API", ".env, Vercel 환경변수"),
        ("NEXT_PUBLIC_SUPABASE_ANON_KEY", "supabase.com → Settings → API", ".env, Vercel 환경변수"),
        ("OPENROUTER_API_KEY", "openrouter.ai → Keys", ".env, Vercel 환경변수"),
        ("(Google Client ID/Secret)", "console.cloud.google.com", "Supabase Authentication 설정 화면"),
    ]
    top = Inches(2.1)
    row_h = Inches(0.7)
    col_w = [Inches(4.3), Inches(4.0), Inches(3.5)]
    x0 = Inches(0.7)
    for r, row in enumerate(env_rows):
        x = x0
        is_head = r == 0
        for c, text in enumerate(row):
            w = col_w[c]
            fill = CORAL if is_head else (WHITE if r % 2 else CREAM2)
            add_rect(s, x, top + row_h * r, w, row_h - Pt(2), fill=fill, line_color=LINE, line_w=0.75)
            add_text(s, text, x + Inches(0.15), top + row_h * r, w - Inches(0.3), row_h - Pt(2),
                      size=13 if not is_head else 14.5, color=WHITE if is_head else INK, bold=is_head,
                      anchor=MSO_ANCHOR.MIDDLE, font=CODE_FONT if not is_head else FONT)
            x += w
    add_text(s, "구글 로그인은 Supabase가 대신 관리하므로, 우리 코드의 .env에는 별도로 등록하지 않습니다.",
              Inches(0.7), Inches(2.1) + row_h * 5 + Inches(0.15), Inches(11.6), Inches(0.5), size=13.5,
              color=MUTED, italic=True)
    set_notes(s, "다섯 사이트를 다 돌아본 뒤, \"결국 무엇을 어디에 붙여넣었는지\"를 한 장으로 정리해주는 "
                 "실용적인 슬라이드다. 실습 중 이 표를 옆에 띄워두고 따라 하게 해도 좋다.")

    s = content_slide(prs, "Part 4 · 계정 준비", "5개 사이트, 한 장으로 정리", nx(), TOTAL)
    table = [
        ("사이트", "한 줄 역할", "비유"),
        ("github.com", "코드 저장·버전 관리, 다른 서비스와의 연결 통로", "사진첩 겸 다리"),
        ("vercel.com", "코드를 실제 인터넷 주소로 배포", "공연장"),
        ("supabase.com", "로그인·데이터베이스·사진 저장소", "창고 겸 대기실"),
        ("openrouter.ai", "여러 회사의 AI 모델을 빌려 호출", "AI 섭외처"),
        ("console.cloud.google.com", "구글 로그인용 열쇠(OAuth) 발급", "신분증 발급소"),
    ]
    col_w = [Inches(3.5), Inches(5.2), Inches(3.1)]
    top = Inches(2.05)
    row_h = Inches(0.72)
    x0 = Inches(0.7)
    for r, row in enumerate(table):
        x = x0
        is_head = r == 0
        for c, text in enumerate(row):
            w = col_w[c]
            fill = CORAL if is_head else (WHITE if r % 2 else CREAM2)
            add_rect(s, x, top + row_h * r, w, row_h - Pt(2), fill=fill, line_color=LINE, line_w=0.75)
            add_text(s, text, x + Inches(0.15), top + row_h * r, w - Inches(0.3), row_h - Pt(2),
                      size=14.5 if not is_head else 15.5, color=WHITE if is_head else INK,
                      bold=is_head, anchor=MSO_ANCHOR.MIDDLE)
            x += w
    set_notes(s, "5개 사이트로 늘어난 요약표. 사진 찍어두라고 권해도 좋다.")

    s = content_slide(prs, "Part 4 · 계정 준비", "다섯 곳이 준비되면, 이렇게 연결됩니다", nx(), TOTAL)
    add_image_fit(s, os.path.join(DIAGRAMS, "architecture.png"), Inches(0.7), Inches(1.9), Inches(11.9),
                  Inches(4.85))
    set_notes(s, "Part3의 아키텍처 그림으로 되돌아와, 방금 하나씩 가입한 다섯 사이트가 그림 속 어느 "
                 "위치에 해당하는지 다시 짚어준다. GitHub은 그림에는 직접 나오지 않지만, Vercel이 코드를 "
                 "가져오는 출발점이라는 것을 설명으로 보충한다.")

    s = content_slide(prs, "Part 4 · 계정 준비", "실습 체크리스트", nx(), TOTAL)
    checks = [
        "GitHub 계정을 만들고 저장소를 하나 만들었다.",
        "Vercel 계정을 만들고 GitHub 저장소를 연결했다.",
        "Supabase 프로젝트를 만들고 데이터베이스 URL·키를 확인했다.",
        "OpenRouter API 키를 발급받았다.",
        "Google Cloud Console에서 OAuth 클라이언트 ID를 발급받았다 (구글 로그인을 쓸 경우).",
    ]
    for i, c in enumerate(checks):
        y = Inches(2.0) + Inches(0.8) * i
        box = slide_box = add_rect(s, Inches(0.7), y, Inches(0.5), Inches(0.5), fill=None,
                                    line_color=CORAL_DARK, line_w=2, radius=0.15)
        add_text(s, c, Inches(1.4), y, Inches(11.0), Inches(0.5), size=17, color=INK,
                  anchor=MSO_ANCHOR.MIDDLE)
    set_notes(s, "이 체크리스트를 실습 진행 상황판처럼 활용한다. 다섯 개를 다 마치지 못했더라도 Part 5로 "
                 "넘어가는 데는 문제가 없다고 안심시켜준다.")

    # ================= PART 5: 개발 과정 기술 심화 =================
    section_divider(prs, 5, "개발 과정 기술 심화", "PRD와 계정 준비가 끝났다면, 이제 실제로 만듭니다")
    nx()

    s = code_slide(
        prs, nx(), TOTAL, "Part 5 · 기술 심화", "실제 코드 저장소, 폴더 구조 먼저 훑어보기",
        "찰칵레시피 저장소 최상위 폴더 구조",
        [
            "src/app/          화면(페이지)들 — 업로드, 재료확인,",
            "                  레시피추천, 설정, 평가 화면 등",
            "src/components/   여러 화면에서 재사용하는 부품",
            "                  (버튼, 로고, 스피너 등)",
            "src/lib/providers/  비전·텍스트 AI 연동 코드",
            "                  (mock 구현과 실제 OpenRouter 구현)",
            "src/lib/fridge/   업로드·인식·추천을 처리하는 서버 코드",
            "src/lib/ratings/  AI 모델 품질 평가 집계 로직",
            "supabase/         데이터베이스 스키마(SQL) 파일들",
            "docs/             PRD 문서와 다이어그램",
            "scripts/          이 강의 자료를 만든 파이썬 스크립트",
        ],
        "폴더 이름만 봐도 무엇이 들어있는지 짐작할 수 있도록 짓는 것이 좋은 습관입니다. 오늘 이후 "
        "슬라이드에서 나오는 파일 경로(예: src/lib/providers/mock-vision.ts)는 전부 이 구조 안에 있는 "
        "실제 파일입니다.",
        "본격적으로 코드 스니펫이 등장하기 전에, 학생들이 \"이 파일이 프로젝트의 어디쯤에 있는지\" 감을 "
        "잡을 수 있도록 지도를 먼저 보여준다.",
    )

    s = content_slide(prs, "Part 5 · 기술 심화", "Next.js와 React, 짧게만 알아보기", nx(), TOTAL)
    add_bullets(s, [
        ("React는 화면을 \"버튼\", \"카드\"처럼 작은 부품(컴포넌트) 단위로 조립해 만드는 방식입니다.", 0),
        ("Next.js는 그 React 부품들을 실제 웹사이트로 만들어주는 틀(프레임워크)입니다. 화면 전환, 서버 "
         "코드 실행까지 함께 챙겨줍니다.", 0),
        ("찰칵레시피의 src/app 폴더 안 폴더 하나하나가 그대로 화면 주소 하나가 됩니다. 예를 들어 "
         "src/app/upload 폴더는 실제 /upload 주소 화면이 됩니다.", 0),
        ("오늘 우리는 이 규칙들을 외울 필요는 없습니다. 클로드 코드가 Next.js의 규칙에 맞게 파일을 "
         "대신 만들어주기 때문입니다. 다만 \"폴더 이름이 곧 화면 주소\"라는 감각만 있으면 코드가 훨씬 "
         "친숙하게 느껴집니다.", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.4), size=17.5)
    set_notes(s, "React/Next.js를 깊게 가르치는 자리가 아니므로 딱 필요한 만큼만 짚어준다. \"폴더 구조가 "
                 "곧 화면 주소\"라는 것만 기억해도 이후 코드 스니펫들을 훨씬 편하게 따라올 수 있다.")

    s = code_slide(
        prs, nx(), TOTAL, "Part 5 · 기술 심화", "서버 액션(Server Action)이란?",
        "개념 예시 (실제 코드를 단순화한 형태)",
        [
            "// 사진 업로드 버튼을 눌렀을 때 서버에서",
            "// 실행되는 함수, 이것이 \"서버 액션\"이다",
            "async function uploadPhoto(formData) {",
            "  \"use server\";  // 이 줄이 서버에서 실행됨을 표시",
            "",
            "  const file = formData.get(\"photo\");",
            "  // 1) Supabase Storage에 사진 저장",
            "  // 2) 비전 API 호출해 재료 인식",
            "  // 3) 인식 결과를 데이터베이스에 저장",
            "  return { success: true };",
            "}",
        ],
        "브라우저 화면의 버튼과, 서버에서 실제로 실행되는 함수를 별도의 API 주소 없이 직접 연결해주는 "
        "Next.js의 기능입니다. \"use server\"라는 한 줄만 붙이면, 이 함수는 사용자의 브라우저가 아니라 "
        "Vercel의 서버 쪽에서 안전하게 실행됩니다. API 키처럼 민감한 값은 이렇게 서버에서만 실행되는 "
        "코드 안에 둬야 안전합니다.",
        "서버 액션은 이 프로젝트 전반에서 가장 자주 쓰이는 패턴이다. \"버튼 클릭이 서버 코드를 직접 "
        "부른다\"는 감각만 잡아줘도 충분하다. 지나치게 기술적으로 파고들 필요는 없다.",
    )

    s = content_slide(prs, "Part 5 · 기술 심화", "이정표 — 계획했던 것과 실제로 만든 것", nx(), TOTAL)
    add_image_fit(s, os.path.join(DIAGRAMS, "milestones.png"), Inches(0.7), Inches(1.95), Inches(11.9),
                  Inches(4.8))
    set_notes(s, "처음 8단계 계획(초록 카드)은 전부 완료됐고, 그 뒤로 계획에 없던 확장 작업(보라 카드)이 "
                 "6개나 더 생겼다는 것을 짚어준다.")

    s = process_slide(
        prs, nx(), TOTAL, "1", "mock(가짜)으로 먼저 전체 그림을 완성합니다.",
        "진짜 AI API를 연결하기 전에, 가짜 응답을 반환하는 \"모형\"으로 화면 전체 흐름부터 완성했습니다.",
        [
            ("이유: API 키가 아직 없어도, 화면과 흐름이 맞는지부터 검증할 수 있습니다.", 0),
            ("사진 업로드 → (가짜) 재료 인식 → (가짜) 레시피 추천까지 끝까지 눌러볼 수 있는 \"뼈대\"를 먼저 "
             "완성했습니다.", 0),
            ("뼈대가 튼튼하면, 나중에 가짜를 진짜로 바꾸는 일은 훨씬 쉬워집니다.", 0, {"bold": True}),
            ("이 지혜는 코딩을 넘어 모든 일에 통합니다. 완벽한 부품보다, 먼저 전체를 얼기설기 이어보는 "
             "것입니다.", 0),
        ],
        "이 단계가 신입생에게 특히 중요한 교훈이다. \"완벽하게 하나씩 끝내고 다음으로\"가 아니라 "
        "\"어설프더라도 처음부터 끝까지 먼저 이어보기\"라는 전략을 강조한다.",
        image=os.path.join(DIAGRAMS, "architecture_mock.png"),
    )

    s = code_slide(
        prs, nx(), TOTAL, "Part 5 · 기술 심화", "mock 코드, 실제로 이렇게 생겼습니다",
        "src/lib/providers/mock-vision.ts (실제 저장소 파일, 지금도 그대로 남아 있습니다)",
        [
            "// 실제 비전 API 대신, 그럴듯한 재료 세트를",
            "// 돌아가며 반환하는 가짜(mock) 구현",
            "const SAMPLE_SETS: DetectedIngredient[][] = [",
            "  [{ name: \"감자\", quantityText: \"3개\" },",
            "   { name: \"양파\", quantityText: \"1개\" }, ...],",
            "  ...",
            "];",
            "",
            "export const mockVisionProvider: VisionProvider = {",
            "  id: \"mock-vision-basic\",",
            "  async detectIngredients(imageUrls) {",
            "    await sleep(500); // 실제 API처럼 살짝 기다린다",
            "    return SAMPLE_SETS[seed % SAMPLE_SETS.length];",
            "  },",
            "};",
        ],
        "실제 사진을 전혀 보지 않고, 미리 준비한 3세트의 \"그럴듯한 재료 목록\" 중 하나를 돌려가며 "
        "돌려줍니다. 이 정도만으로도 업로드→인식→추천까지 전체 화면 흐름을 완성하고 검증할 수 있습니다. "
        "이 파일은 지금도 src/lib/providers/mock-vision.ts 경로에 그대로 남아 있고, 저장소 GitHub 페이지에서 "
        "직접 열어볼 수 있습니다.",
        "실제 코드를 그대로 보여주는 것이 핵심이다. \"mock이 뭔가 대단히 복잡한 시뮬레이션\"이라는 오해를 "
        "깨고, 사실은 몇 줄짜리 하드코딩이라는 걸 확인시켜준다. 학생들에게 \"이 정도는 나도 짤 수 있겠다\"는 "
        "감각을 준다.",
    )

    s = process_slide(
        prs, nx(), TOTAL, "2", "진짜 AI로 교체합니다.",
        "가짜 응답을 지우지 않고 그대로 둔 채, openrouter.ai의 실제 비전·텍스트 모델을 새로 연결했습니다.",
        [
            ("mock 코드는 지우지 않고 남겨뒀습니다. \"가짜 로직 vs 실제 API\"를 비교하는 수업 자료로도 "
             "쓰기 위해서입니다. (src/lib/providers/mock-vision.ts, mock-text.ts)", 0),
            ("한 모델이 한도 초과나 오류로 응답을 못 하면, 다음 모델로 자동으로 갈아타는 안전장치를 "
             "추가했습니다.", 0),
            ("실제 AI 응답은 형식이 들쭉날쭉할 수 있어, 응답을 관대하게 해석하는 처리도 더했습니다.", 0),
            ("\"완벽한 AI\"를 기다리지 않고, 여러 겹의 안전망으로 현실적인 신뢰도를 만들었습니다.", 0,
             {"bold": True, "color": CORAL_DARK}),
        ],
        "이 단계에서 \"AI도 가끔 틀리거나 응답을 못 할 수 있다\"는 사실을 솔직하게 알려준다.",
        image=os.path.join(DIAGRAMS, "architecture.png"),
    )

    s = code_slide(
        prs, nx(), TOTAL, "Part 5 · 기술 심화", "실제 모델 후보 목록과 자동 재시도(fallback)",
        "src/lib/providers/openrouter-vision.ts (실제 저장소 파일)",
        [
            "const CANDIDATES = [",
            "  { id: \"openrouter-vision-1\",",
            "    model: \"google/gemma-4-26b-a4b-it:free\" },",
            "  { id: \"openrouter-vision-2\",",
            "    model: \"nvidia/nemotron-nano-12b-v2-vl:free\" },",
            "  { id: \"openrouter-vision-3\",",
            "    model: \"google/gemma-4-31b-it:free\" },",
            "];",
            "",
            "// 고른 모델이 죽어있을 때를 대비해",
            "// openrouter/free(자동 라우팅)를 안전망으로 붙인다",
            "const candidates = [model, \"openrouter/free\"];",
        ],
        "드롭다운에서 사용자가 고른 모델을 1순위로 시도하고, 실패하면 openrouter/free(그 순간 살아있는 "
        "무료 모델을 OpenRouter가 대신 골라주는 라우터)로 자동 전환합니다. 무료 모델은 가끔 응답이 "
        "불안정할 수 있다는 것을 전제로, 이렇게 여러 겹의 안전망을 코드 레벨에서 미리 설계해뒀습니다.",
        "chatJsonWithFallback이라는 공통 함수가 이 candidates 배열을 순서대로 시도한다는 것을 언급해도 "
        "좋다. 실패 처리를 하나의 함수로 통일해뒀다는 것이 코드 품질 측면에서 중요한 포인트다.",
        bullets=[
            ("1순위: 사용자가 화면에서 고른 모델.", 0),
            ("2순위(안전망): openrouter/free 자동 라우팅.", 0),
            ("두 번 다 실패하면 사용자에게 친절한 오류 메시지를 보여줍니다.", 0),
        ],
    )

    s = content_slide(prs, "Part 5 · 기술 심화", "실제로 등록된 무료 모델 목록", nx(), TOTAL)
    model_rows = [
        ("구분", "모델 ID", "비고"),
        ("비전 1순위", "google/gemma-4-26b-a4b-it:free", "기본 선택"),
        ("비전 2순위", "nvidia/nemotron-nano-12b-v2-vl:free", "작은 모델, 반복 응답 방어 로직 포함"),
        ("비전 3순위", "google/gemma-4-31b-it:free", ""),
        ("비전 자동", "openrouter/free", "안전망 겸 자동 선택 옵션"),
        ("텍스트 1순위", "nvidia/nemotron-3-super-120b-a12b:free", "기본 선택"),
        ("텍스트 2순위", "qwen/qwen3-next-80b-a3b-instruct:free", ""),
        ("텍스트 3순위", "cognitivecomputations/dolphin-mistral-24b-venice-edition:free", ""),
        ("AI 판정용", "openrouter/free", "평가 대상과 겹치지 않는 모델로 고정"),
    ]
    top = Inches(1.95)
    row_h = Inches(0.53)
    col_w = [Inches(1.7), Inches(6.9), Inches(3.2)]
    x0 = Inches(0.7)
    for r, row in enumerate(model_rows):
        x = x0
        is_head = r == 0
        for c, text in enumerate(row):
            w = col_w[c]
            fill = CORAL if is_head else (WHITE if r % 2 else CREAM2)
            add_rect(s, x, top + row_h * r, w, row_h - Pt(2), fill=fill, line_color=LINE, line_w=0.75)
            add_text(s, text, x + Inches(0.15), top + row_h * r, w - Inches(0.3), row_h - Pt(2),
                      size=12.5 if not is_head else 14, color=WHITE if is_head else INK, bold=is_head,
                      anchor=MSO_ANCHOR.MIDDLE, font=CODE_FONT if (c == 1 and not is_head) else FONT)
            x += w
    add_text(s, "2026년 7월 기준 실제 등록 목록입니다. OpenRouter의 무료 모델 라인업은 자주 바뀌므로, "
              "직접 만들 때는 openrouter.ai에서 price = free 필터로 최신 목록을 다시 확인해야 합니다.",
              Inches(0.7), top + row_h * len(model_rows) + Inches(0.15), Inches(11.8), Inches(0.6),
              size=12.5, color=MUTED, italic=True)
    set_notes(s, "실제로 docs/PRD.md 7.1에 정리된 표를 그대로 옮긴 것이다. 무료 모델 목록이 자주 바뀐다는 "
                 "점을 반드시 언급해, 학생들이 \"이 모델 이름을 외워야 한다\"고 오해하지 않게 한다.")

    s = code_slide(
        prs, nx(), TOTAL, "Part 5 · 기술 심화", "보안 규칙(RLS)도 실제로 이렇게 코드로 작성합니다",
        "supabase/09_public_recipe_ratings.sql (실제 저장소 파일)",
        [
            "-- 로그인한 회원이면 누구나 레시피를 볼 수 있게",
            "create policy \"recipes_authenticated_select\"",
            "  on recipes",
            "  for select",
            "  using (auth.uid() is not null);",
            "",
            "-- 레시피 평가(좋아요/싫어요/코멘트)도 마찬가지로",
            "create policy \"recipe_feedback_authenticated_select\"",
            "  on recipe_feedback",
            "  for select",
            "  using (auth.uid() is not null);",
        ],
        "\"auth.uid() is not null\"은 \"로그인한 사용자라면\"이라는 뜻입니다. 이 규칙 덕분에 레시피와 "
        "그 평가는 회원이면 누구나 볼 수 있게 됐지만, 냉장고 사진과 인식된 재료는 이 규칙을 적용하지 "
        "않아 여전히 본인만 볼 수 있습니다.",
        "\"사례 8\"(관리자 전용 → 공개 전환)에서 실제로 실행했던 SQL 그대로다. RLS 정책이 실제로는 몇 "
        "줄짜리 SQL 문장이라는 걸 보여줘, \"보안 규칙\"이라는 말이 주는 위압감을 낮춰준다.",
    )

    s = process_slide(
        prs, nx(), TOTAL, "3", "배포하고, 실제로 써보고, 고칩니다.",
        "Vercel에 배포한 뒤에도 끝이 아닙니다. 실제 화면에서 눌러보며 계속 다듬었습니다.",
        [
            ("GitHub에 코드를 올리면 Vercel이 자동으로 새 버전을 배포합니다.", 0),
            ("서버리스 함수(실제 API를 호출하는 서버 코드)는 기본 실행 제한 시간이 있어서, 응답이 오래 "
             "걸릴 수 있는 AI 호출 구간은 maxDuration 설정으로 시간을 늘려줬습니다.", 0),
            ("배포 후 실제 사진을 올려 \"비전 API → 재료 인식 → 텍스트 API → 레시피 생성\"이 진짜로 "
             "도는지 직접 확인했습니다.", 0),
            ("\"다 만들었다\"가 아니라 \"써보니 이건 고쳐야겠다\"의 반복이 진짜 개발입니다.", 0,
             {"bold": True}),
        ],
        "실제로 이번 강의 자료를 준비하며 배포된 사이트에 직접 접속해 사진을 올려 테스트했던 경험을 "
        "생생하게 들려준다.",
    )

    s = process_slide(
        prs, nx(), TOTAL, "4", "대화하다 보면, 계획에 없던 것도 생깁니다.",
        "PRD의 8단계 계획을 다 마친 뒤에도, 대화를 계속하며 6개의 새로운 작업이 자연스럽게 생겨났습니다.",
        [
            ("\"이해하기 어려운 앱 이름을 바꿔야겠다.\" → 브랜드 리뉴얼(찰칵레시피 이름과 로고)로 "
             "확장됐습니다.", 0),
            ("\"다른 회원 레시피도 볼 수 있으면 좋겠다.\" → 관리자 전용 화면을 전체 공개 구조로 "
             "재설계했습니다.", 0),
            ("\"AI 모델끼리 성능을 비교해볼 수 있을까?\" → AI 모델 품질 평가 시스템을 아예 새로 "
             "설계했습니다.", 0),
            ("계획은 나침반이지, 감옥이 아닙니다. 대화가 곧 기획의 연장입니다.", 0,
             {"bold": True, "color": CORAL_DARK}),
        ],
        "이 슬라이드는 \"바이브 코딩은 계획을 안 세운다\"는 오해를 바로잡는다. v1에서 이 문장이 \"관리자만 "
        "보이는 이름을 바꿔야겠다\"였는데, 실제로는 처음부터 recipe4fridge_pic이라는 이름 자체가 "
        "누구에게나 낯설고 이해하기 어려웠다는 문제였다. \"이해하기 어려운 앱 이름\"으로 바로잡았다.",
    )

    s = content_slide(prs, "Part 5 · 기술 심화", "미니 사례 — AI가 AI를 평가하게 만들다", nx(), TOTAL)
    add_image_fit(s, os.path.join(DIAGRAMS, "model_eval.png"), Inches(0.7), Inches(1.85), Inches(11.9),
                  Inches(4.9))
    set_notes(s, "이 사례는 오늘 강의에서 가장 고급스러운 예시다. 사용자가 1~5점으로 평가하는 것에 더해, "
                 "AI 모델 하나를 \"채점관\"으로 세워 다른 모델들의 결과를 자동으로 채점하게 만들었다. AI "
                 "판정 상자의 부제(\"인식·추천할 때마다 서버가 자동 채점\")를 손가락으로 짚으며, 이게 매번 "
                 "자동으로 일어난다는 것을 분명히 설명한다.")

    s = content_slide(prs, "Part 5 · 기술 심화", "데이터 모델도 결국 상식입니다", nx(), TOTAL)
    add_image_fit(s, os.path.join(DIAGRAMS, "erd.png"), Inches(0.7), Inches(1.85), Inches(11.9),
                  Inches(4.9))
    set_notes(s, "복잡해 보이는 데이터베이스 설계도(ERD)도, 실은 우리가 일상에서 쓰는 \"표\"들의 관계일 "
                 "뿐이라는 걸 강조한다. 화면 좌상단의 범례(화살표 방향의 의미)를 먼저 짚어준 뒤, 회원 한 "
                 "명이 여러 세션을 갖고, 세션 하나가 여러 사진과 재료를 갖고, 그 결과로 레시피가 생기는 "
                 "흐름을 화살표를 손가락으로 짚어가며 설명하면 좋다.")

    # ================= PART 6: 실사용 피드백과 대응 사례 연구 =================
    section_divider(prs, 6, "테스트하고 보완한 이야기", "배포는 끝이 아니라 시작입니다 — 실제 사용자 피드백과 대응")
    nx()

    s = content_slide(prs, "Part 6 · 테스트와 보완", "이 파트에서 볼 것 — 진짜 있었던 일들", nx(), TOTAL)
    add_bullets(s, [
        ("지금부터 보여드릴 사례는 전부 이 강의를 준비하며 실제로 있었던 요청과 대응입니다. 지어낸 예시가 "
         "아닙니다.", 0, {"bold": True, "color": CORAL_DARK}),
        ("각 사례는 \"실제 요청 → 진단 → 해결 → 배운 점\" 네 단계로 정리했습니다.", 0),
        ("배운 점은 기술, 디자인, 제품, 방법론 네 가지 관점으로 나눠 표시했습니다. 버그 하나를 고치는 "
         "일도 여러 관점에서 배울 게 있다는 것을 보여드리고 싶었습니다.", 0),
        ("화면 오른쪽의 스크린샷 자리는 실제 배포된 화면을 캡처해 넣는 자리입니다. 강의 자료를 준비하며 "
         "직접 채워 넣으시면 됩니다.", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.4), size=18)
    set_notes(s, "이 파트가 이번 v2.0에서 가장 힘을 준 신규 파트라는 것을 알린다. 스크린샷 placeholder에 "
                 "대해서도 미리 설명해, 강의자가 captured/ 폴더에 실제 캡처를 채워 넣으면 완성된다는 것을 "
                 "알려준다.")

    s = case_study_slide(
        prs, nx(), TOTAL, 1,
        "재료 인식 결과를 보여줄 때, 사진이 더 크면 좋겠어. 확대 기능을 넣어도 좋고.",
        "기존 화면은 업로드한 사진을 180×180 크기의 작은 썸네일로만 보여주고 있었습니다. 인식이 맞는지 "
        "확인하려면 사진 속 글씨나 재료가 잘 안 보였습니다.",
        "PhotoGallery라는 새 컴포넌트를 만들어, 썸네일을 240×240으로 키우고 클릭하면 화면 전체 크기의 "
        "어두운 오버레이 위에 원본 사진을 크게 띄우도록 했습니다. 클릭하면 다시 닫힙니다.",
        "\"작은 썸네일 하나\"가 아니라 \"작은 미리보기 + 필요할 때 확대\"라는 2단계 정보 구조가 훨씬 "
        "친절합니다.", "디자인",
        "case01_photo_zoom.png",
        "재료 확인 화면에서 업로드한 사진을 클릭해 전체화면으로 확대한 상태를 캡처하세요.",
        "이 사례는 순수하게 UX 개선 요청이다. 기술적으로는 어렵지 않지만, \"정보를 언제 작게, 언제 크게 "
        "보여줄지\"를 결정하는 디자인 감각을 보여주는 좋은 예시다.",
    )

    s = case_study_slide(
        prs, nx(), TOTAL, 2,
        "레시피 추천 화면에서 모델 드롭다운 UI가 카드 영역을 벗어나.",
        "긴 OpenRouter 모델 이름이 드롭다운 안에 들어가면서, 부모 요소인 flex 컨테이너의 기본값 "
        "min-width: auto 때문에 드롭다운이 카드 폭보다 넓게 밀려나며 카드 밖으로 삐져나왔습니다.",
        "해당 select 요소에 min-width: 0과 max-width: 100%를 주고, 텍스트 모델 select에는 flex-basis: "
        "100%를 추가해 아예 한 줄을 통째로 차지하도록 CSS를 수정했습니다.",
        "flex 레이아웃에서 자식 요소는 기본적으로 \"내용물보다 작아지지 않으려는\" 성질이 있습니다. 이걸 "
        "모르면 이런 오버플로 버그를 화면을 직접 눌러보기 전까지 발견하기 어렵습니다.", "기술",
        "case02_dropdown_fix.png",
        "재료 확인 화면 하단, 텍스트 API 선택 드롭다운이 카드 안에 깔끔하게 들어간 모습을 캡처하세요.",
        "CSS flexbox의 min-width: auto 기본값은 프론트엔드 개발에서 매우 흔한 함정이다. 이 사례를 통해 "
        "\"눈으로 봐야만 발견되는 버그\"가 있다는 것을 알려준다.",
    )

    s = case_study_slide(
        prs, nx(), TOTAL, 3,
        "AI 모델 평가 화면에서 표시되는 모델의 이름과 재료 인식·레시피 추천 화면에 표시되는 모델의 "
        "이름이 달라서 혼란스러워.",
        "점수판 화면은 데이터베이스에 저장된 provider_id(예: openrouter-vision-1)를 그대로 보여주고 "
        "있었는데, 업로드·추천 화면은 사람이 읽는 라벨(예: OpenRouter · google/gemma...)을 보여주고 "
        "있었습니다. 같은 모델이 서로 다른 이름으로 두 곳에 나타난 것입니다.", "점수판 컴포넌트에 라벨 "
        "사전(visionProviders, textProviders)을 전달해, provider_id를 사람이 읽는 라벨로 변환해 "
        "보여주도록 통일했습니다.",
        "같은 개념은 앱 전체에서 반드시 같은 이름으로 불러야 합니다. 내부 식별자(id)와 사용자에게 보여줄 "
        "이름(label)을 처음부터 구분해 설계해두면 이런 문제를 예방할 수 있습니다.", "디자인",
        "case03_model_name_match.png",
        "AI 모델 평가(/models) 화면에서, 업로드 화면과 동일한 모델명이 표시되는 점수판을 캡처하세요.",
        "일관성(consistency)이라는 디자인 원칙을 실제 버그 사례로 설명할 수 있는 좋은 기회다.",
    )

    s = case_study_slide(
        prs, nx(), TOTAL, 4,
        "레시피 추천을 요청했다가 에러가 났는데, 다시 성공했는데도 이전 에러 메시지가 화면에 그대로 "
        "남아 있어.",
        "브라우저의 \"뒤로 가기\" 캐시(bfcache) 때문에, 페이지가 새로 불러와지는 대신 이전 상태(에러 "
        "메시지가 떠 있던 상태)를 그대로 복원하고 있었습니다. 코드상의 로직 오류가 아니라 브라우저의 "
        "동작 방식 때문에 생긴 버그였습니다.",
        "pageshow라는 브라우저 이벤트를 감지해, event.persisted가 true일 때(캐시에서 복원됐을 때) "
        "에러·상태 값을 초기화하도록 처리했습니다.",
        "겉보기엔 코드 버그처럼 보여도, 실제로는 브라우저 자체의 동작 방식을 알아야 원인을 찾을 수 있는 "
        "문제들이 있습니다. 재현이 잘 안 되는 버그일수록 \"뒤로 가기\", \"새로고침\" 같은 브라우저 동작을 "
        "의심해봐야 합니다.", "기술",
        "case04_stale_error_fixed.png",
        "에러 발생 후 성공적으로 재요청해, 이전 에러 메시지 없이 깔끔한 화면을 캡처하세요.",
        "bfcache는 신입생에게 생소한 개념이니 \"페이지를 다시 그리지 않고 그대로 꺼내오는 브라우저의 "
        "지름길 기능\" 정도로 쉽게 풀어준다.",
    )

    s = case_study_slide(
        prs, nx(), TOTAL, 5,
        "레시피 저장 버튼을 눌러도 아이콘이 그대로인 것 같아. 확인해줘.",
        "코드를 다시 살펴봤지만 저장 버튼 아이콘을 바꾸는 로직 자체는 문제가 없었습니다. 이미 배포된 "
        "코드에 반영되어 있었습니다. 코드가 아니라 재현 여부부터 확인이 필요한 상황이었습니다.",
        "짐작으로 코드를 고치는 대신, 배포된 사이트에 직접 접속해 실제로 저장 버튼을 눌러봤습니다. "
        "클릭하자 아이콘이 회색 테두리에서 채워진 주황 원으로 정상적으로 바뀌었고, \"저장한 레시피\" "
        "목록에도 정확히 반영되는 것을 확인했습니다. 결국 사용자의 브라우저 캐시 문제였을 가능성이 "
        "높다고 결론지었습니다.",
        "버그 신고를 받으면 코드부터 의심하기 전에, 먼저 재현해보는 것이 순서입니다. 재현이 안 되면 "
        "브라우저 캐시나 사용자 환경 문제일 수 있습니다. \"고쳤다\"고 말하기 전에 실제로 확인하는 습관이 "
        "중요합니다.", "방법론",
        "case05_save_icon.png",
        "레시피 저장 버튼을 클릭해, 채워진 주황색 아이콘으로 바뀐 순간을 캡처하세요.",
        "이 사례는 이번 강의 준비 중 실제로 라이브 브라우저 테스트로 검증했던 과정 그대로다. \"확인 없이 "
        "고쳤다고 말하지 않는다\"는 태도를 강조한다.",
    )

    s = case_study_slide(
        prs, nx(), TOTAL, 6,
        "bakery 테마에서 버튼에 마우스를 올렸을 때 색이 이상해 보여.",
        "세 가지 디자인 테마(apricot, greens, bakery) 중 bakery 테마만 hover(마우스를 올렸을 때) 색상 "
        "값이 다른 테마와 다른 방식으로 정의되어 있었습니다. 테마마다 CSS 변수를 따로 관리하다 보니 "
        "생긴 불일치였습니다.",
        "세 테마의 hover 색상 정의 방식을 통일하고, bakery 테마의 값을 다른 두 테마와 같은 규칙으로 "
        "다시 맞췄습니다.",
        "디자인 시스템에서 \"같은 역할을 하는 값\"은 테마가 여러 개라도 반드시 같은 규칙으로 정의해야 "
        "합니다. 그렇지 않으면 테마를 하나 늘릴 때마다 이런 불일치가 계속 생깁니다.", "디자인",
        "case06_bakery_theme.png",
        "설정 화면에서 bakery 테마를 선택한 뒤, 버튼에 마우스를 올린 hover 상태를 캡처하세요.",
        "디자인 시스템(테마 토큰)을 잘 관리하지 않으면 생기는 전형적인 문제다. CSS 변수로 색상을 "
        "관리하는 이유를 설명하기 좋은 사례다.",
    )

    s = case_study_slide(
        prs, nx(), TOTAL, 7,
        "recipe4fridge_pic은 앱의 기술적 이름으로만 쓰면 좋겠어. 이해하기 어려운 앱 이름을 바꿔야겠다.",
        "저장소 이름이자 초기 서비스명이었던 recipe4fridge_pic이 사용자에게 그대로 노출되고 있었는데, "
        "영어 약자라 직관적으로 와닿지 않는다는 문제가 있었습니다. 넓은 화면과 좁은 화면 양쪽에 이름이 "
        "중복 노출되는 문제도 있었습니다.",
        "\"찰칵레시피\"라는 새 브랜드명과 카메라 모양의 로고(BrandMark 컴포넌트)를 만들고, "
        "recipe4fridge_pic은 저장소·패키지명 같은 기술적 용도로만 남겼습니다. 브랜드는 넓은 화면에서는 "
        "좌측 사이드바 상단에, 좁은 화면에서는 상단바에 한 곳에서만 보이도록 정리했습니다.",
        "제품 이름은 개발자가 아니라 사용할 사람의 눈높이에서 지어야 합니다. 기술적인 이름과 사용자에게 "
        "보여줄 이름을 분리하는 것만으로도 첫인상이 크게 달라집니다.", "제품",
        "case07_brand_rename.png",
        "홈 화면 상단에 \"찰칵레시피\" 로고와 이름이 표시된 모습을 캡처하세요.",
        "브랜딩은 종종 사소한 요청처럼 보이지만, 사용자가 서비스를 처음 만나는 순간의 인상을 결정한다는 "
        "점에서 제품 전략의 일부다.",
    )

    s = case_study_slide(
        prs, nx(), TOTAL, 8,
        "다른 회원의 레시피에 대한 평가는 회원들의 인기 투표 개념으로 접근하고, 관리자만 보는 게 아니라 "
        "일반 회원들도 집계 결과를 확인할 수 있도록 해줘.",
        "기존에는 관리자(profiles.is_admin = true)만 볼 수 있는 별도 대시보드에서 레시피 평가를 집계해 "
        "보고 있었습니다. 이 구조로는 일반 회원들끼리 서로의 레시피를 보고 평가하는 \"인기 투표\"가 "
        "불가능했습니다.",
        "Supabase의 보안 규칙(RLS)을 \"본인 데이터만\"에서 \"로그인한 회원이면 누구나 레시피와 그 평가를 "
        "볼 수 있다\"로 완화하고, 관리자 전용 대시보드를 삭제한 뒤 회원 공개용 화면(/ratings)을 새로 "
        "만들었습니다. 단, 냉장고 사진과 인식된 재료는 여전히 본인만 볼 수 있도록 남겨뒀습니다.",
        "권한(누가 무엇을 볼 수 있는가)을 바꾸는 요청은 화면 하나만 고치는 일이 아닙니다. 데이터베이스 "
        "보안 규칙부터 다시 설계해야 하는, 파급 범위가 큰 변경이라는 것을 이해해야 합니다.", "기술",
        "case08_public_ratings.png",
        "레시피 평가(/ratings) 화면에서 다른 회원의 레시피에 좋아요·코멘트를 남긴 화면을 캡처하세요.",
        "\"권한 모델\"이라는 개념을 실제 사례로 소개하기 좋다. 화면 하나를 공개하는 것이 데이터베이스 "
        "정책 전체를 다시 설계하는 일로 이어진다는 파급효과를 보여준다.",
    )

    s = content_slide(prs, "Part 6 · 테스트와 보완", "사례 9 — AI 모델 평가, 어떻게 설계할지 함께 정했습니다", nx(), TOTAL)
    add_bullets(s, [
        ("\"AI 모델 품질을 사용자와 AI가 함께 평가하게 만들고 싶다\"는 아이디어는 정할 것이 많았습니다.", 0),
        ("그래서 클로드 코드가 선택지를 정리해 질문으로 제시했고, 사람이 하나씩 결정했습니다.", 0,
         {"bold": True, "color": CORAL_DARK}),
        ("자동화 수준: 점수판만 보여줄지, 점수에 따라 드롭다운 순서까지 자동으로 바꿀지 → \"완전 "
         "자동화\"를 선택했습니다.", 0),
        ("판정 시점: 사용자가 요청할 때만 채점할지, 매번 자동으로 채점할지 → \"매번 자동\"을 "
         "선택했습니다.", 0),
        ("판정에 쓸 모델: Claude 같은 유료 API를 쓸지, 무료 모델을 쓸지 → 비용을 이유로 openrouter/free를 "
         "선택했습니다.", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.5), size=17)
    set_notes(s, "이 사례는 스크린샷보다 \"의사결정 과정\" 자체가 핵심이라 카드형이 아니라 목록형으로 "
                 "구성했다. AI가 답을 혼자 정하는 게 아니라, 트레이드오프를 정리해 사람에게 선택지로 "
                 "제시하고 사람이 결정한다는 협업 방식을 보여준다. 이는 Part 8의 \"성찰을 통한 목표 설정\", "
                 "\"클로드 코드 협업에 대한 신뢰\" 역량과 바로 연결된다.")

    s = case_study_slide(
        prs, nx(), TOTAL, 10,
        "배포는 했는데, 실제로 이 링크로 접속해서 사진을 올려보면 진짜로 동작하는지 확인이 필요해.",
        "코드가 배포됐다고 해서 실제로 잘 동작한다는 보장은 없습니다. 실제 API 키, 실제 네트워크 환경, "
        "실제 브라우저에서 확인해야 합니다.",
        "배포된 사이트에 직접 로그인해, 실제 냉장고 사진을 업로드했습니다. 비전 API가 실제로 재료를 "
        "인식하고(계란, 양배추, 오이, 토마토 등), 텍스트 API가 그 재료로 레시피를 생성하는 전체 흐름을 "
        "눈으로 확인했습니다.",
        "\"코드가 배포됐다\"와 \"실제로 동작한다\"는 다른 이야기입니다. 배포 후 실사용 시나리오로 직접 "
        "검증하는 단계를 절대 생략해서는 안 됩니다.", "방법론",
        "case10_live_test.png",
        "실제 냉장고 사진을 업로드해, 비전 API가 인식한 재료와 텍스트 API가 만든 레시피가 함께 보이는 "
        "화면을 캡처하세요.",
        "이 사례는 이번 강의를 준비하며 실제로 진행했던 D4(배포 후 실제 API 검증) 작업 그대로다. QA "
        "(품질보증)라는 개념을 자연스럽게 소개할 수 있다.",
    )

    s = case_study_slide(
        prs, nx(), TOTAL, 11,
        "다른 레시피 더 보기를 눌렀는데, 방금 받은 것과 너무 비슷한 레시피가 또 나와.",
        "재요청할 때 이전에 추천받은 레시피 목록을 프롬프트에 함께 알려주지 않아서, 텍스트 AI가 매번 "
        "비슷한 답을 반복하는 경향이 있었습니다.",
        "재요청 시 이전 추천 이력을 프롬프트에 포함해 \"이전과 겹치지 않게 추천해달라\"고 명시적으로 "
        "지시하도록 바꿨습니다. 화면에도 이전 추천 이력을 함께 누적해서 보여주도록 했습니다.",
        "AI에게 맥락(이전에 무엇을 이미 보여줬는지)을 충분히 알려주지 않으면, 사람 입장에서는 뻔한 것도 "
        "AI는 계속 반복할 수 있습니다. 필요한 맥락을 매번 프롬프트에 실어 보내는 습관이 중요합니다.",
        "기술",
        "case11_recipe_variety.png",
        "레시피 추천 화면에서 이번 추천과 이전 추천 이력이 함께 표시된 모습을 캡처하세요.",
        "\"AI는 대화 맥락을 스스로 기억하지 않는다\"는 점을 짚어주기 좋은 사례다. 매 요청마다 필요한 "
        "정보를 사람이 다시 실어 보내야 한다는 것을 강조한다.",
    )

    s = case_study_slide(
        prs, nx(), TOTAL, 12,
        "예전에 찍었던 냉장고 사진 기록도 삭제할 수 있으면 좋겠어.",
        "세션(사진 업로드 한 번의 기록)을 저장만 할 뿐, 삭제할 방법이 화면에 없었습니다. 잘못 올렸거나 "
        "필요 없어진 기록을 치울 수 없었습니다.",
        "이전 재료 목록 화면에 세션 삭제 버튼을 추가하고, 세션을 지우면 연결된 사진·인식된 재료·추천 "
        "레시피까지 데이터베이스에서 함께 정리되도록(on delete cascade) 만들었습니다.",
        "데이터를 지울 때는 그 데이터에 연결된 다른 데이터까지 함께 고려해야 합니다. 하나만 지우고 "
        "나머지가 고아처럼 남으면 데이터베이스가 지저분해집니다.", "기술",
        "case12_session_delete.png",
        "이전 재료 목록 화면에서 세션 삭제 버튼과 확인 과정을 캡처하세요.",
        "on delete cascade라는 데이터베이스 개념을 이 사례로 소개할 수 있다. \"연결된 것까지 한 번에 "
        "정리한다\"는 뜻이라고 쉽게 풀어준다.",
    )

    s = content_slide(prs, "Part 6 · 테스트와 보완", "열두 가지 사례에서 배운 것을 정리하면", nx(), TOTAL)
    add_bullets(s, [
        ("기술 관점: CSS의 숨은 기본값, 브라우저 캐시처럼 눈에 안 보이는 동작 방식을 알아야 진짜 원인을 "
         "찾을 수 있습니다.", 0),
        ("디자인 관점: 작은 것 하나(썸네일 크기, 모델 이름 표기)의 일관성이 사용자 경험 전체를 좌우합니다.", 0),
        ("제품 관점: 이름 하나, 권한 하나를 바꾸는 결정도 사용자와 서비스 전체에 영향을 줍니다.", 0),
        ("방법론 관점: 짐작으로 고치지 않고, 재현하고 검증하는 습관이 신뢰할 수 있는 결과를 만듭니다.", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=18)
    set_notes(s, "Part6을 마무리하며 네 가지 관점을 한 번 더 정리해준다. 이 네 관점이 Part 8의 7가지 "
                 "역량과 자연스럽게 이어진다는 것을 예고해도 좋다.")

    # ================= PART 7 =================
    section_divider(prs, 7, "여러분도 할 수 있다는 확신", "\"나는 코드를 몰라서\"는 이제 핑계가 안 됩니다")
    nx()

    s = content_slide(prs, "Part 7 · 자신감", "\"나는 코드를 몰라서 못한다\"는 착각", nx(), TOTAL)
    add_bullets(s, [
        ("오늘 다룬 화면 흐름, 아키텍처, 데이터 모델, 다섯 사이트 연결까지 전부 문장으로 설명해서 "
         "만들었습니다.", 0),
        ("문법을 몰라도, 무엇을 원하는지를 구체적으로 말할 수 있으면 시작할 수 있습니다.", 0),
        ("처음엔 완벽한 문장이 아니어도 됩니다. 대화를 주고받으며 점점 정확해지면 됩니다.", 0),
        ("진짜 필요한 것은 문법 지식이 아니라, 좋은 질문을 던지는 습관입니다.", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "이 슬라이드에서 신입생들의 가장 큰 심리적 장벽(\"나는 전공자가 아니라서\")을 정면으로 "
                 "다룬다.")

    s = content_slide(prs, "Part 7 · 자신감", "오늘 앱도 결국 대화의 누적이었습니다", nx(), TOTAL)
    add_bullets(s, [
        ("\"냉장고 사진으로 레시피 추천하는 앱 만들고 싶어.\" 씨앗이 된 첫 문장이었습니다.", 0),
        ("\"사진이 더 크면 좋겠어, 확대 기능을 넣어도 좋고.\" 써보고 나서 나온 개선 요청이었습니다.", 0),
        ("\"이 판정도 무료 모델을 쓰면 어떨까?\" 비용을 고민하다 나온 아이디어였습니다.", 0),
        ("작은 문장 수백 개가 쌓여서, 지금의 찰칵레시피가 되었습니다.", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "실제로 이 강의 자료를 준비하며 오갔던 대화 문장들을 그대로 인용한 것임을 밝히면 "
                 "설득력이 커진다.")

    s = content_slide(prs, "Part 7 · 자신감", "여러분이라면 무엇을 만들겠습니까?", nx(), TOTAL)
    add_bullets(s, [
        ("동아리 출석과 회비를 정리해주는 앱을 만들 수 있습니다.", 0),
        ("자취생을 위한 \"이번 주 장보기 리스트 추천\" 앱을 만들 수 있습니다.", 0),
        ("전공 수업 조모임 일정을 자동으로 맞춰주는 앱을 만들 수 있습니다.", 0),
        ("내 강아지·고양이 사진을 보고 오늘 컨디션을 기록해주는 앱을 만들 수 있습니다.", 0),
        ("이 목록의 공통점은, 전부 내 주변의 작은 불편에서 출발한다는 것입니다.", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "이 슬라이드는 참여형으로 진행한다. 예시를 던진 뒤 \"여러분 주변엔 어떤 불편이 있나요?\" "
                 "라고 실제로 몇 명에게 질문해 짧게 답을 들어본다.")

    s = content_slide(prs, "Part 7 · 자신감", "실패해도 괜찮은 이유", nx(), TOTAL)
    add_bullets(s, [
        ("모든 변경 이력은 GitHub에 기록되어 남습니다. 문제가 생기면 이전 상태로 되돌릴 수 있습니다.", 0),
        ("오늘 배운 것처럼 mock(가짜)으로 먼저 시도해보고, 안전하게 진짜로 바꿔갈 수 있습니다.", 0),
        ("AI에게 \"이거 왜 이렇게 됐어?\"라고 물어보면, 원인을 함께 찾아줍니다.", 0),
        ("완벽한 한 번이 아니라, 작게 자주 시도하고 되돌릴 수 있는 여러 번이 이 시대의 안전망입니다.", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "실패에 대한 두려움을 낮추는 슬라이드다. Part 4에서 배운 버전 관리(GitHub 커밋)를 다시 "
                 "연결지어 설명하면 좋다.")

    # ================= PART 8 =================
    section_divider(prs, 8, "AI 활용 개발자에게 필요한 7가지 역량", "코딩 실력보다 오래가는 것")
    nx()

    s = content_slide(prs, "Part 8 · 마무리", "왜 \"코딩 실력\"만으로는 부족해지는가?", nx(), TOTAL)
    add_bullets(s, [
        ("문법을 외우고 빠르게 타이핑하는 능력은, AI가 이미 사람보다 잘합니다.", 0),
        ("그렇다면 이제 사람에게 남는 것, 오히려 더 중요해지는 것은 무엇일까요?", 0),
        ("오늘 찰칵레시피를 만드는 과정에서 실제로 계속 쓰였던 능력 7가지를 정리했습니다.", 0,
         {"bold": True, "color": CORAL_DARK}),
        ("이 7가지는 전공과 무관하게, 지금 이 자리의 모든 신입생에게 이미 씨앗이 있는 능력들입니다.", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "Part8 전체의 문을 여는 슬라이드다. \"코딩을 잘하는 사람\"에서 \"AI와 함께 좋은 것을 "
                 "만드는 사람\"으로 기준이 이동했다는 것을 분명히 선언한다.")

    competencies = [
        ("사람들의 욕구를 파악해서\n꿈을 키워내는 능력",
         "기술이 아니라 사람에서 출발합니다.",
         [
             "\"냉장고 사진 → 오늘 뭐 해먹지\"는 기술이 아니라 사람의 불편에서 시작됐습니다.",
             "좋은 개발자는 코드보다 먼저, 사람이 진짜 원하는 것을 알아채는 사람입니다.",
             "작은 불편을 발견하는 눈이, 곧 좋은 아이디어의 씨앗입니다.",
         ]),
        ("성찰을 통한\n위치 파악 및 목표 설정 능력",
         "지금 어디에 있고, 어디로 가는지 알아야 합니다.",
         [
             "PRD를 먼저 쓴 것도 결국 우리가 지금 뭘 하려는지를 스스로 확인하는 과정이었습니다.",
             "이정표(마일스톤)를 계속 되돌아보며 계획대로인지, 방향을 바꿔야 하는지 점검했습니다.",
             "AI는 방향을 대신 정해주지 않습니다. 방향은 내가 정하고, AI는 그 길을 함께 걷습니다.",
         ]),
        ("상식적 판단력",
         "\"이게 말이 되나?\"를 끝까지 놓지 않는 힘입니다.",
         [
             "AI가 내놓은 결과를 \"그럴듯해 보인다\"로 끝내지 않고 정말 맞는지 검증했습니다.",
             "평가 점수가 5건도 안 쌓였는데 순위를 확정 짓지 않은 것도 상식적 판단입니다.",
             "AI의 속도에 상식이라는 브레이크를 다는 사람이, 신뢰할 수 있는 결과를 만듭니다.",
         ]),
        ("클로드 코드 협업에 대한 신뢰",
         "동료를 믿고 일을 맡길 줄 아는 능력입니다.",
         [
             "모든 걸 직접 타이핑하려 하지 않고, 방향을 정한 뒤 실행은 믿고 맡겼습니다.",
             "결과를 함께 검토하고, 아니다 싶으면 다시 요청하는 건강한 협업을 반복했습니다.",
             "신뢰는 맹신이 아닙니다. 확인하고, 고치고, 다시 신뢰하는 과정의 반복입니다.",
         ]),
        ("개선에 대한\n열망과 열정",
         "\"이 정도면 됐다\"에서 한 걸음 더 나아갑니다.",
         [
             "사진이 작아 보이면 확대 기능을, 이름이 밋밋하면 브랜드를 새로 고민했습니다.",
             "완성됐다고 끝내지 않고, 써보면서 계속 더 나아질 부분을 찾아냈습니다.",
             "이 열정이 있어야, 계획에 없던 6개의 확장 작업도 마다하지 않고 해낼 수 있었습니다.",
         ]),
        ("포기하지 않는\n의지와 체력",
         "안 되면 될 때까지, 다르게 다시 시도하는 힘입니다.",
         [
             "무료 AI 모델이 오류를 내면, 될 때까지 다른 모델로 자동 전환하는 안전망을 만들었습니다.",
             "화면이 깨지고, 배포가 실패하고, 파일 권한이 꼬여도 원인을 찾아 끝까지 고쳤습니다.",
             "긴 작업일수록 오늘은 여기까지를 알고, 다시 이어갈 체력을 관리하는 것도 실력입니다.",
         ]),
        ("차별적 성과 창출",
         "남들과 같은 것을 만들지 않는 용기입니다.",
         [
             "인기 투표와 AI 모델 품질 평가를 분리한 설계처럼, 남이 안 하는 디테일까지 챙겼습니다.",
             "\"AI가 AI를 평가하게 만드는\" 시스템처럼, 한 걸음 더 나아간 아이디어를 실제로 구현했습니다.",
             "결국 도구는 누구에게나 열려 있습니다. 차이를 만드는 건 그 도구로 무엇을 어디까지 밀어붙이는가 "
             "입니다.",
         ]),
    ]
    for idx, (title, essence, bullets) in enumerate(competencies, start=1):
        s = competency_slide(prs, nx(), TOTAL, idx, title, essence, bullets,
                              f"역량 {idx} 상세 설명. 반드시 오늘 강의에서 실제로 등장했던 구체적 장면과 "
                              f"연결지어 이야기한다. 추상적 훈화로 끝나지 않도록, \"오늘 우리가 이걸 이렇게 "
                              f"써먹었죠\"라는 방식으로 되짚어준다.")

    s = content_slide(prs, "Part 8 · 마무리", "7가지 역량, 한 장으로 정리", nx(), TOTAL)
    labels = [
        "① 욕구 파악 · 꿈을 키우는 능력", "② 성찰을 통한 위치 파악 · 목표 설정",
        "③ 상식적 판단력", "④ 클로드 코드 협업에 대한 신뢰",
        "⑤ 개선에 대한 열망 · 열정", "⑥ 포기하지 않는 의지 · 체력",
        "⑦ 차별적 성과 창출",
    ]
    cols = 2
    cw, ch = Inches(5.75), Inches(0.85)
    gx, gy = Inches(0.3), Inches(0.2)
    for i, label in enumerate(labels):
        r, c = divmod(i, cols)
        x = Inches(0.7) + c * (cw + gx)
        y = Inches(2.0) + r * (ch + gy)
        add_rect(s, x, y, cw, ch, fill=CREAM2, radius=0.15)
        add_text(s, label, x + Inches(0.25), y, cw - Inches(0.5), ch, size=17, color=INK, bold=True,
                  anchor=MSO_ANCHOR.MIDDLE)
    set_notes(s, "전체 요약 슬라이드다. 이 슬라이드를 사진 찍어가라고 권해도 좋다.")

    s = content_slide(prs, "Part 8 · 마무리", "오늘의 작은 과제", nx(), TOTAL)
    add_bullets(s, [
        ("내 주변의 작은 불편 하나를 문장으로 적어봅니다. (예: OOO 하기가 매번 귀찮다.)", 0),
        ("그 불편을 풀어주는 앱을 한 문단으로 설명해봅니다. 이것이 나만의 첫 PRD가 됩니다.", 0),
        ("오늘 배운 내용대로, GitHub·Vercel·Supabase·OpenRouter를 준비해 나만의 앱을 완성해봅니다.", 0),
        ("완벽하지 않아도 됩니다. 오늘 배운 4단계 루프처럼, 다음 시간에 함께 다듬어가면 됩니다.", 0,
         {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=19)
    set_notes(s, "다음 수업과의 연결고리다. v1에서는 \"한 곳에 가입만이라도 해보기\"였는데, 오늘 v2.0은 "
                 "PRD 작성부터 5개 사이트 실습, 기술 심화, 테스트 사례까지 훨씬 실전적으로 다뤘으므로 "
                 "\"나만의 앱을 완성해보기\"로 목표 수준을 높였다.")

    s = content_slide(prs, "부록", "오늘 강의에서 다루지 않은 것들", nx(), TOTAL)
    add_bullets(s, [
        ("정확한 칼로리 계산이나 영양 정보 제공은 다루지 않았습니다. 처음 PRD의 비목표에 명시해뒀던 "
         "범위입니다.", 0),
        ("재고 관리(유통기한 알림 등)나 모바일 네이티브 앱도 이번 프로젝트의 범위 밖입니다.", 0),
        ("사용자가 아주 많아졌을 때의 대규모 트래픽 대응, 유료 결제 연동 같은 주제도 오늘은 다루지 "
         "않았습니다.", 0),
        ("무엇을 하지 않을지 정하는 것도 PRD의 중요한 역할이라는 것을, 여기서 다시 한번 확인할 수 "
         "있습니다.", 0, {"bold": True, "color": CORAL_DARK}),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=18)
    set_notes(s, "강의의 범위를 명확히 하는 슬라이드다. Part3에서 배운 \"비목표\"라는 개념을 다시 "
                 "떠올리게 하며 마무리로 이어간다. 학생들이 \"오늘 배운 게 전부가 아니다\"라는 것과 "
                 "\"오늘 배운 원칙은 그런 더 큰 주제에도 그대로 적용된다\"는 것을 함께 알려준다.")

    s = content_slide(prs, "부록", "자주 묻는 질문", nx(), TOTAL)
    faqs = [
        ("무료로 계속 쓸 수 있나요?", "다섯 사이트 모두 개인 프로젝트 수준에서는 무료 요금제로 충분합니다. "
         "다만 사용자가 많아지면 유료 전환이 필요해질 수 있습니다."),
        ("클로드 코드가 없으면 못 만드나요?", "아닙니다. 다른 AI 코딩 도구로도 비슷한 방식이 가능합니다. "
         "다만 오늘 배운 \"설명 → 구현 → 검증 → 반복\"이라는 사고방식 자체가 핵심입니다."),
        ("한 번에 완벽하게 만들어야 하나요?", "전혀 아닙니다. 오늘 본 것처럼 mock으로 시작해서 여러 번 "
         "고쳐나가는 것이 정상적인 과정입니다."),
        ("프로그래밍을 하나도 몰라도 되나요?", "오늘 수준의 결과물은 몰라도 시작할 수 있습니다. 다만 화면 "
         "구조, 데이터 흐름 같은 기본 개념을 알수록 AI에게 더 정확하게 설명할 수 있습니다."),
    ]
    top = Inches(2.0)
    row_h = Inches(1.05)
    for i, (q, a) in enumerate(faqs):
        y = top + row_h * i
        add_rect(s, Inches(0.7), y, Inches(11.8), row_h - Pt(8), fill=CREAM2, radius=0.1)
        add_text(s, f"Q. {q}", Inches(0.95), y + Pt(6), Inches(11.3), Inches(0.35), size=15.5, color=CORAL_DARK,
                  bold=True)
        add_text(s, f"A. {a}", Inches(0.95), y + Inches(0.42), Inches(11.3), Inches(0.55), size=13.5,
                  color=INK, line_spacing=1.2)
    set_notes(s, "강의를 마치기 전 학생들이 흔히 품는 의문을 미리 해소해준다. 실제 질의응답 시간에 나온 "
                 "질문이 있다면 이 슬라이드에 계속 추가해나가도 좋다.")

    s = new_slide(prs, bg=CREAM)
    add_rect(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill=CORAL)
    add_text(s, "마치며", Inches(0.9), Inches(1.5), Inches(10), Inches(0.5), size=18, color=CORAL_DARK,
              bold=True)
    add_text(s, "무엇이든 만들 수 있습니다,\n여러분이 정말 원한다면.", Inches(0.9), Inches(2.1),
              Inches(11.5), Inches(1.9), size=40, color=INK, bold=True, line_spacing=1.25)
    add_text(s,
             "저는 40년을 코드와 함께 살았지만, 오늘 여러분과 함께 본 이 여정이야말로 지금 이 시대에 "
             "코딩을 배운다는 것의 진짜 의미라고 생각합니다. 냉장고 사진 한 장에서 시작한 작은 아이디어가, "
             "실제로 접속할 수 있는 서비스가 되기까지, 문법이 아니라 대화가 그 다리를 놓았습니다.\n\n"
             "여러분에게도 이미 그 씨앗이 있습니다. 오늘 배운 것을 가지고, 여러분만의 첫 문장을 "
             "시작해보시기 바랍니다.",
             Inches(0.9), Inches(4.15), Inches(11.5), Inches(2.6), size=18, color=INK, line_spacing=1.4)
    set_notes(s, "마지막 인사다. 진심을 담아 천천히 읽어준다.")

    n_final = nx()
    s = content_slide(prs, "부록", "참고 링크와 자료 모음", n_final, TOTAL)
    add_bullets(s, [
        ("배포된 앱: recipe4fridge-pic.vercel.app", 0),
        ("GitHub 저장소: 오늘 본 모든 코드, 커밋 이력, mock 구현 파일이 그대로 있습니다.", 0),
        ("PRD 문서: docs/PRD.md, 그리고 단계별 스냅샷 docs/PRD00~PRD06.", 0),
        ("오늘 본 다이어그램: docs/diagrams/ (아키텍처, 이정표, AI 평가 구조도, ERD, 전체 여정 지도).", 0),
        ("5개 사이트: github.com · vercel.com · supabase.com · openrouter.ai · console.cloud.google.com.", 0),
    ], Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.3), size=18)
    set_notes(s, "강의가 끝난 뒤에도 학생들이 스스로 다시 찾아볼 수 있는 참고 자료 목록이다. 실제 URL과 "
                 "경로로 교체해 배포하면 좋다.")

    return prs, TOTAL, n_final


if __name__ == "__main__":
    prs, total, n = build()
    print(f"planned TOTAL={total}, n counter={n}, actual slide count={len(prs.slides)}")
    out_path = os.path.join(PROJECT_ROOT, "docs", "찰칵레시피_바이브코딩_강의_v2_완전판.pptx")
    prs.save(out_path)
    print("wrote", out_path)


